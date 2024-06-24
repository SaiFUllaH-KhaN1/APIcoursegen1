from langchain_community.chat_models import ChatOpenAI
# from langchain_openai import ChatOpenAI
from langchain.chains import LLMChain
from langchain.prompts import BaseChatPromptTemplate, PromptTemplate
import base64
from langchain.prompts import PromptTemplate
from langchain.schema.messages import HumanMessage, SystemMessage
# from PyPDF2 import PdfReader
from pypdf import PdfReader
from langchain_community.vectorstores import FAISS
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter, RecursiveCharacterTextSplitter
import os
import shutil
import json
from langchain.schema.document import Document
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders import UnstructuredExcelLoader, UnstructuredWordDocumentLoader, UnstructuredPowerPointLoader, TextLoader
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
from langchain_core.pydantic_v1 import BaseModel, Field
from typing import List, Dict, Any, Optional
from docx2python import docx2python
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import shutil
import re
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from PIL import Image
from langchain.utils.math import cosine_similarity


def RAG(file_content,embeddings,file,session_var):
    print("file is:",file)
    
    filename = file.filename
    print(filename)
    extension = filename.rsplit('.', 1)[1].lower()
    filename_without_extension = filename.rsplit('.', 1)[0].lower()
    output_path_byfile = f"./imagefolder_{session_var}/images_{session_var}_{filename_without_extension}"
    if not os.path.exists(output_path_byfile):
        os.makedirs(output_path_byfile) 
    print("Extension is:",extension)
    raw_text = ''
    texts = '' # for pdf image path only!
    if extension=="pdf":

        # For Image processing
        if f'extracted_content{session_var}.pdf' not in filename:
            pdf_reader = PdfReader(file_content)
            pgcount=0
            for i, page in enumerate(pdf_reader.pages):
                text_instant = page.extract_text()
                pgcount += 1
                imgcount = 1
                text_instant = f"\nThe Content of PageNumber:{pgcount} of file name:{filename_without_extension} is:\n{text_instant}.\nEnd of PageNumber:{pgcount} of file name:{filename_without_extension}\n"
                try:
                    for image_file_object in page.images:
                        # base_name = os.path.splitext(os.path.basename(image_file_object.name))[0]  # Get the base file name without extension
                        extension = os.path.splitext(image_file_object.name)[1]  # Get the file extension
                        base_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}"  # Construct new file name with count
                        if extension == ".jp2":
                            print("Checking Image Extension",extension)
                            image_path = os.path.join(output_path_byfile, base_name)  # Construct the full output path
                            imgcount += 1
                            with open(image_path, "wb") as fp:
                                fp.write(image_file_object.data)
                            with Image.open(image_path) as im:
                                new_image_name = base_name + ".png"  # New image name for PNG
                                new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
                                im.save(new_image_path)  # Save the image as PNG
                                os.remove(image_path)
                        else:
                            image_name = f"FileName {filename_without_extension} PageNumber {pgcount} ImageNumber {imgcount}{extension}"  # Construct new file name with count
                            image_path = os.path.join(output_path_byfile, image_name)  # Construct the full output path
                            imgcount += 1
                            with open(image_path, "wb") as fp:
                                fp.write(image_file_object.data)
                            with Image.open(image_path) as im:
                                if im.mode in ["P", "PA"]:
                                    print("Image of P or PA Mode detected:",im.mode)
                                    im = im.convert("RGBA")  # Convert palette-based images to RGBA
                                    new_image_name = base_name + ".png"  # New image name for PNG
                                    new_image_path = os.path.join(output_path_byfile, new_image_name)  # New full path for PNG
                                    im.save(new_image_path)  # Save the image as PNG
                                    os.remove(image_path)
                                else:
                                    pass
                        

                except Exception as e:
                    print(f"Error processing image {image_file_object.name}: {e}")
                    continue  # Skip to the next image                        
                if text_instant:
                    texts += text_instant
            # at this point we got text and images extracted, stored in ./images folder, lets summarize images

            # Get image summaries
            # image_elements = []
            # image_summaries = []
            # image_path = output_path_byfile

            # def encode_image(image_path):
            #     with open(image_path, "rb") as f:
            #         return base64.b64encode(f.read()).decode('utf-8')

            # def summarize_image(encoded_image, basename):
            #     prompt = [
            #         SystemMessage(content="You are a bot that is good at analyzing images."),
            #         HumanMessage(content=[
            #             {
            #                 "type": "text",
            #                 "text": f"Describe the contents of this image. Tell what FileName, PageNumber and ImageNumber of this image is by seeing this information: {basename}. Your output should look like this: 'This image that belongs to FileName: ..., PageNumber: ..., ImageNumber: .... In this Image ...'"
            #             },
            #             {
            #                 "type": "image_url",
            #                 "image_url": {
            #                     "url": f"data:image/jpeg;base64,{encoded_image}"
            #                 },
            #             },
            #         ])
            #     ]
            #     response = ChatOpenAI(model="gpt-4o", max_tokens=128).invoke(prompt)
            #     return response.content

            # for i in os.listdir(output_path_byfile):
            #     if i.endswith(('.png', '.jpg', '.jpeg')):
            #         image_path = os.path.join(output_path_byfile, i)
            #         basename = os.path.basename(image_path) 
            #         print(os.path.basename(image_path))
            #         encoded_image = encode_image(image_path)
            #         image_elements.append(encoded_image)
            #         summary = summarize_image(encoded_image,basename)
            #         image_summaries.append(summary)

            # print("Image Summary is::",image_summaries)

            # text_merged = texts + "\n" + str(image_summaries)
            # print(text_merged) # Will be used onwards to create the faiss_index text database
            
        else:
            print("We are processing for url or youtube's extracted_content.pdf")
            # Without the Image Processing
            doc_reader = PdfReader(file_content)

            for i, page in enumerate(doc_reader.pages):
                text = page.extract_text()

                if text:
                    raw_text += text

    elif extension=="csv":
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        file.save(temp_path)
        loader = CSVLoader(file_path=temp_path)
        data = loader.load()
        raw_text = raw_text.join(document.page_content + '\n\n' for document in data)
        os.remove(temp_path)

    elif extension=="xlsx" or extension=="xls":
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        file.save(temp_path)
        loader = UnstructuredExcelLoader(temp_path)
        data = loader.load()
        raw_text = raw_text.join(document.page_content + '\n\n' for document in data)
        os.remove(temp_path)

    elif extension=="docx":
        # temp_path = os.path.join(filename)
        # file.seek(0)
        # file.save(temp_path)
        # loader = UnstructuredWordDocumentLoader(temp_path)
        # data = loader.load()
        # raw_text = raw_text.join(document.page_content for document in data)
        # os.remove(temp_path)
        def extract_and_rename_images(docx_path, output_dir):
            # Extract the contents of the DOCX file
            content = docx2python(docx_path, extract_image=True,  image_folder=output_dir)
            print(content)
            # Flatten the list of lists containing the images
            images = content.images
            image_info = []
            for image_name in images.keys():
                base_name, ext = os.path.splitext(image_name)
                print(f"Base Name: {base_name}, Extension: {ext}")
                image_info.append((base_name, ext))
            # Rename images based on their location
            image_count = 1
            i = 0
            texts = ""
            for section_index, section in enumerate(content.body):
                for row_index, row in enumerate(section):
                    for cell_index, cell in enumerate(row):
                        for paragraph_index, paragraph in enumerate(cell):
                            if '----media/' in paragraph:
                                # Extract the image filename from the placeholder
                                start_index = paragraph.find('----media/') + len('----media/')
                                end_index = paragraph.find('----', start_index)
                                image_filename = paragraph[start_index:end_index]

                                # Construct the original image path
                                original_image_path = os.path.join(output_dir, image_filename)
                                # base_name = os.path.splitext(os.path.basename(image_file_object.name))[0]  # Get the base file name without extension
                                # extension = os.path.splitext(image_file_object.name)[1]  # Get the file extension
                                # Create a new image name based on its location
                                base_name,ext = image_info[i]
                                new_image_name = f"FileName {filename_without_extension} PageNumber Null ImageNumber {image_count}{ext}"
                                new_image_path = os.path.join(output_dir, new_image_name)

                                # Rename the image
                                if os.path.exists(original_image_path):
                                    shutil.move(original_image_path, new_image_path)

                                    # Update the placeholder in the paragraph
                                    new_placeholder = f"----media/ImageNumber:{image_count} PageNumber:Null of FileName:{filename_without_extension}----"
                                    paragraph = paragraph.replace(f"----media/{image_filename}----", new_placeholder)
                                    cell[paragraph_index] = paragraph
                                    image_count += 1
                                    i += 1
                            # print(paragraph)
                            texts += paragraph + "\n"
            return texts

        texts = extract_and_rename_images(file_content, output_path_byfile)
        print("texts is:::",texts)

    elif extension=="pptx":
        # temp_path = os.path.join(filename)
        # file.seek(0)
        # file.save(temp_path)
        # loader = UnstructuredPowerPointLoader(temp_path)
        # data = loader.load()
        # raw_text = raw_text.join(document.page_content for document in data)
        # os.remove(temp_path)
        def iter_picture_shapes(prs):
            slide_number = 1
            # image_number = 1
            # img_names = []
            for slide in prs.slides:
                image_number = 1
                print("slide",slide)
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for s in shape.shapes:
                            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image = shape.image
                                print("image_number",image_number)
                                image_filename = f'FileName {filename_without_extension} SlideNumber {slide_number} ImageNumber {image_number}.{image.ext}'
                                # img = f'SlideNumber:{slide_number} of FileName:{filename_without_extension}-ImageNumber {image_number}'
                                image_number += 1
                                print(image_filename)
                                # img_names.append(img)

                                image_path = os.path.join(output_path_byfile, image_filename)
                                with open(image_path, "wb") as fp:
                                    fp.write(image.blob)

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        print("image_number",image_number)
                        image_filename = f'FileName {filename_without_extension} SlideNumber {slide_number} ImageNumber {image_number}.{image.ext}'
                        # img = f'SlideNumber:{slide_number} of FileName:{filename_without_extension} with ImageNumber:{image_number}\n'
                        image_number += 1
                        print(image_filename)
                        # img_names.append(img)
                        image_path = os.path.join(output_path_byfile, image_filename)
                        with open(image_path, "wb") as fp:
                            fp.write(image.blob)
                slide_number += 1  


        iter_picture_shapes(Presentation(file_content))

        # langchain unstructuredworddoc method
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        file.save(temp_path)
        loader = UnstructuredPowerPointLoader(temp_path,mode='elements')
        data = loader.load()
        print("data",data)

        # Step 1: Collect content for each page number
        page_contents = {}
        for doc in data:
            # Access the metadata and page content correctly
            page_number = doc.metadata.get('page_number')
            if page_number is not None:
                if page_number not in page_contents:
                    page_contents[page_number] = []
                page_contents[page_number].append(doc.page_content)

        # Step 2: Combine the content for each page number
        # combined_page_contents = [{'page_number': page,'filename': filename, 'page_content': ' '.join(contents)} for page, contents in page_contents.items()]

        combined_page_contents = [
            {
                'SlideNumber': page,
                'FileName': filename_without_extension,
                'page_content': f"{' '.join(contents)} End of SlideNumber:{page} with Filename:{filename_without_extension} ----"
            }
            for page, contents in page_contents.items()
        ]

        texts = str(combined_page_contents)
        print(texts)
        os.remove(temp_path)

    elif extension=="txt":
        print(f"TExt file name is ::{filename}")
        temp_path = os.path.join(f"{filename}{session_var}")
        file.seek(0)
        file.save(temp_path)
        loader = TextLoader(temp_path)
        data = loader.load()
        raw_text = raw_text.join(document.page_content for document in data)
        os.remove(temp_path)

    # chunking recursively without semantic search, this does not uses openai embeddings for chunking
    text_splitter = RecursiveCharacterTextSplitter(
    chunk_size = 1536,
    chunk_overlap  = 0,
    length_function = len,
    )
    # uses openai embeddings for chunking, costs time and money but gives good performances, not ideal for real-time
    # text_splitter = SemanticChunker(OpenAIEmbeddings(), breakpoint_threshold_type="percentile", number_of_chunks= 10000)
    print("Before Embeddings!")
    
    print("Now Doing Embeddings!")
    
    if texts:
        print("Running Text Merged (Pdfs other than extractedcontent.pdf)")
        text_splitter_formerged = RecursiveCharacterTextSplitter(chunk_size=1536, chunk_overlap=128, length_function=len,)
        text = text_splitter_formerged.split_text(texts)
        docsearch = FAISS.from_texts(text, embeddings)

    elif raw_text:
        raw_text_splitted = text_splitter.split_text(raw_text)
        docsearch = FAISS.from_texts(raw_text_splitted, embeddings)

    print("docsearch made")
    return docsearch

promptSelector = PromptTemplate(
    input_variables=["input_documents","human_input"],
    template="""
    As an educational chatbot, you are tasked with guiding the selection of the most suitable learning scenario 
    tailored to the specific requirements of course content.
    Your decision-making process is 
    informed by evaluating 'Human Input' and 'Input Documents', allowing you to determine the best fit among 
    the following for course development:

    Gamified: A gamified environment that encourages applying subject knowledge to escape a scenario like an Exit Game is designed, 
    enhancing investigative and critical thinking skills.
    Linear: Straightforward, step-by-step training on a topic, ending with quizzes to evaluate understanding.
    Branched: A sandbox-style experience where users can explore various aspects of a topic at 
    their own pace, including subtopics with quizzes. These Byte-size subtopics help in learning being more digestible.
    Simulation: A decision-making driven simulation learning experience, where different choices lead to different 
    outcomes, encouraging exploration of pertinent consequences faced. Hence, learning is achieved via a simulated experience. 

    'Human Input': ({human_input})
    'Input Documents': ({input_documents})

    Your reply should be one of the below (Depends on what you find most suitable to be selected):
    Bot: Gamified Scenario
    Bot: Simulation Scenario
    Bot: Linear Scenario
    Bot: Branched Scenario
    """
)


prompt_linear = PromptTemplate(
    input_variables=["input_documents","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging educational content in a Linear Scenario Format using
    a system of blocks. You give step-by-step detail information such that you are teaching a student.

    ***WHAT TO DO***
    To accomplish educational Linear Scenario creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted in Linear Scenario structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Linear Scenarios are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'SelfAssessmentTextBlock' with title, and descritpion(This is part of formative assessment. It is assessment of oneself or one's actions, attitudes, or performance in relation to learning objectives.) 
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Linear Scenario: A type of educational structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed information to users based on "Learning Objectives", "Content Areas" and "Input Documents". The use of TextBlocks and MediaBlocks actually act as segregating various aspects of the subject matter, by giving information of the various concepts of subject matter in detailed and dedicated way. For each of the concept or aspect of the subject, a detailed information, illustrative elaboration (if needed) and Question are asked for testing. At the end of covering all aspects of the subject, there will be FeedbackAndFeedforwardBlock and SelfAssessmentTextBlock followed by the TestBlocks having series or single QuestionBlock/s to test user's knowledge and GoalBlock for scoring users.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!   
    
    \nOverview structure of the Linear Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the scenario and proceedings)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give illustrated, complimentary material to elaborate on the information given in Text Blocks. Generate a MediaBlock/s to complement the information provided in Text Blocks. Firstly, see if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then use your imagination to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    QuestionBlock/s (Students after a certain important TextBlock/s or MediaBlock/s are tested via QuestionBlock/s if they learned from the content of the specific block to which this Question Block belongs to. Give atleast 5 QuestionBlocks and so the previous TextBlocks should have enough content to be covered in these 5 QuestionBlocks named as QB1,QB2 till QB5. It can be even higher depending on the course content.)
    FeedbackAndFeedforwardBlock
    SelfAssessmentTextBlock
    GoalBlock
    \nEnd of Overview structure\n

    Problem to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively. 


    \n\nEXAMPLE START: LINEAR SCENARIO:\n\n
{{
      "title": "(Insert a fitting Title Here)",
      "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here) and so on"
        }},
        {{
          "id": "B3",
          "Purpose": "This MANDATORY block (In terms of either one Text Block or multiple per scenario.) is where you !Begin by giving welcome message to the scenario. In further Text Blocks down the example format you use these blocks to give detailed information on every aspect of various subject matters as asked.",
          "type": "TextBlock",
          "title": "(Insert Text Here)",
          "description": "(Insert Text Here)"
        }},
        {{
          "id": "B4",
          "Purpose": "This OPTIONAL block (In terms of either one Media Block or multiple or no Media Block per scenario. In case of no Media Block, Text Block use is Mandatory to give information about each and every aspect of the subject matter) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
          "type": "MediaBlock",
          "title": "(Insert Text Here)",
          "mediaType": "Image(Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
          "description": "(Insert Text Here)",
          "overlayTags": [
            "(Insert Text Here)",
            "(Insert Text Here)"
          ]
        }},
        {{
          "id": "B5",
          "type": "TextBlock",
          "title": "Feedback_And_Feedforward",
          "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
          "id": "B6",
          "type": "TextBlock",
          "title": "Self_Assessment",
          "description": "Self Assessment=(Insert Text Here)"
        }},
        {{
          "id": "QB1",
          "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of this specific branch in regards to its information given in its TextBlocks and MediBlocks. The QuestionBlocks can be single or multiple depending on the content and importance at hand",
          "type": "QuestionBlock",
          "questionText": "(Insert Text Here)",
          "answers": [
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)"
          ],
          "correctAnswer": "(Insert Text Here)",
          "wrongAnswerMessage": "(Insert Text Here)"
        }},
        {{
          "id": "GB",
          "type": "GoalBlock",
          "title": "Congratulations!",
          "score": 3
        }}
      ],
      "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
          "source": "B1",
          "target": "B2"
        }},
        {{
          "source": "B2",
          "target": "B3"
        }},
        {{
          "source": "B3",
          "target": "B4"
        }},
        {{
          "source": "B4",
          "target": "B5"
        }},
        {{
          "source": "B5",
          "target": "B6"
        }},
        {{
          "source": "B6",
          "target": "QB1"
        }},
        {{
          "source": "QB1",
          "target": "GB"
        }}
    ]
}}
    \n\nEND OF EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.

    Chatbot (Tone of a teacher teaching student in great detail):"""
)

prompt_linear_retry = PromptTemplate(
    input_variables=["incomplete_response"],
    template="""
    Based on the INSTRUCTIONS below, an 'Incomplete Response' was created. Your task is to complete
    this response by continuing from exactly where the 'Incomplete Response' discontinued its response.
    Complete the response by continuing exactly from the discontinued point, which is specified by '[CONTINUE_EXACTLY_FROM_HERE]'.
    Never include [CONTINUE_EXACTLY_FROM_HERE] in your response. This is just for your information.
    DO NOT RESPOND FROM THE START OF THE 'Incomplete Response'. Just start from the exact point where the 'Incomplete Response' is discontinued! 
    Take great care into the ID heirarchy considerations while continuing the incomplete response.
    'Incomplete Response': {incomplete_response};

    !!!WARNING: KEEP YOUR RESPONSE SHORT, since you have alreay reached your token limit!!! 

    !!!NOTE: YOU HAVE TO ENCLOSE THE JSON PARENTHESIS BY KEEPING THE 'Incomplete Response' IN CONTEXT!!!
    
    !!!CAUTION: INCLUDE WITH NODES, ALSO RELATIVE EDGES FOR DEFINING CONNECTIONS OF BLOCKS!!!

    BELOW IS THE INSTRUCTION SET BASED ON WHICH THE 'Incomplete Response' WAS CREATED ORIGINALLY:
    INSTRUCTION SET:
    [
    You are an educational bot that creates engaging educational content in a Linear Scenario Format using
    a system of blocks. You give step-by-step detail information such that you are teaching a student.

    ***WHAT TO DO***
    To accomplish educational Linear Scenario creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted in Linear Scenario structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Linear Scenarios are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'SelfAssessmentTextBlock' with title, and descritpion(This is part of formative assessment. It is assessment of oneself or one's actions, attitudes, or performance in relation to learning objectives.) 
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Linear Scenario: A type of educational structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed information to users based on "Learning Objectives", "Content Areas" and "Input Documents". The use of TextBlocks and MediaBlocks actually act as segregating various aspects of the subject matter, by giving information of the various concepts of subject matter in detailed and dedicated way. For each of the concept or aspect of the subject, a detailed information, illustrative elaboration (if needed) and Question are asked for testing. At the end of covering all aspects of the subject, there will be FeedbackAndFeedforwardBlock and SelfAssessmentTextBlock followed by the TestBlocks having series or single QuestionBlock/s to test user's knowledge and GoalBlock for scoring users.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!   
    
    \nOverview structure of the Linear Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the scenario and proceedings)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give illustrated, complimentary material to elaborate on the information given in Text Blocks. Generate a MediaBlock/s to complement the information provided in Text Blocks. Firstly, see if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then use your imagination to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    QuestionBlock/s (Students after a certain important TextBlock/s or MediaBlock/s are tested via QuestionBlock/s if they learned from the content of the specific block to which this Question Block belongs to. Give atleast 5 QuestionBlocks and so the previous TextBlocks should have enough content to be covered in these 5 QuestionBlocks named as QB1,QB2 till QB5. It can be even higher depending on the course content.)
    FeedbackAndFeedforwardBlock
    SelfAssessmentTextBlock
    GoalBlock
    \nEnd of Overview structure\n

    Problem to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively. 


    \n\nEXAMPLE START: LINEAR SCENARIO:\n\n
{{
      "title": "(Insert a fitting Title Here)",
      "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here) and so on"
        }},
        {{
          "id": "B3",
          "Purpose": "This MANDATORY block (In terms of either one Text Block or multiple per scenario.) is where you !Begin by giving welcome message to the scenario. In further Text Blocks down the example format you use these blocks to give detailed information on every aspect of various subject matters as asked.",
          "type": "TextBlock",
          "title": "(Insert Text Here)",
          "description": "(Insert Text Here)"
        }},
        {{
          "id": "B4",
          "Purpose": "This OPTIONAL block (In terms of either one Media Block or multiple or no Media Block per scenario. In case of no Media Block, Text Block use is Mandatory to give information about each and every aspect of the subject matter) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
          "type": "MediaBlock",
          "title": "(Insert Text Here)",
          "mediaType": "Image(Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
          "description": "(Insert Text Here)",
          "overlayTags": [
            "(Insert Text Here)",
            "(Insert Text Here)"
          ]
        }},
        {{
          "id": "B5",
          "type": "TextBlock",
          "title": "Feedback_And_Feedforward",
          "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
          "id": "B6",
          "type": "TextBlock",
          "title": "Self_Assessment",
          "description": "Self Assessment=(Insert Text Here)"
        }},
        {{
          "id": "QB1",
          "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of this specific branch in regards to its information given in its TextBlocks and MediBlocks. The QuestionBlocks can be single or multiple depending on the content and importance at hand",
          "type": "QuestionBlock",
          "questionText": "(Insert Text Here)",
          "answers": [
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)"
          ],
          "correctAnswer": "(Insert Text Here)",
          "wrongAnswerMessage": "(Insert Text Here)"
        }},
        {{
          "id": "GB",
          "type": "GoalBlock",
          "title": "Congratulations!",
          "score": 3
        }}
      ],
      "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
          "source": "B1",
          "target": "B2"
        }},
        {{
          "source": "B2",
          "target": "B3"
        }},
        {{
          "source": "B3",
          "target": "B4"
        }},
        {{
          "source": "B4",
          "target": "B5"
        }},
        {{
          "source": "B5",
          "target": "B6"
        }},
        {{
          "source": "B6",
          "target": "QB1"
        }},
        {{
          "source": "QB1",
          "target": "GB"
        }}
    ]
}}
    \n\nEND OF EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.
    ]


    Chatbot:"""
)

prompt_linear_simplify = PromptTemplate(
    input_variables=["input_documents","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging educational content in a Linear Scenario Format using
    a system of blocks. You give step-by-step detail information such that you are teaching a student.

    ***WHAT TO DO***
    To accomplish educational Linear Scenario creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted in Linear Scenario structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Linear Scenarios are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'SelfAssessmentTextBlock' with title, and descritpion(This is part of formative assessment. It is assessment of oneself or one's actions, attitudes, or performance in relation to learning objectives.) 
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Linear Scenario: A type of educational structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed information to users based on "Learning Objectives", "Content Areas" and "Input Documents". The use of TextBlocks and MediaBlocks actually act as segregating various aspects of the subject matter, by giving information of the various concepts of subject matter in detailed and dedicated way. For each of the concept or aspect of the subject, a detailed information, illustrative elaboration (if needed) and Question are asked for testing. At the end of covering all aspects of the subject, there will be FeedbackAndFeedforwardBlock and SelfAssessmentTextBlock followed by the TestBlocks having series or single QuestionBlock/s to test user's knowledge and GoalBlock for scoring users.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!   
    
    \nOverview structure of the Linear Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the scenario and proceedings)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give illustrated, complimentary material to elaborate on the information given in Text Blocks. Generate a MediaBlock/s to complement the information provided in Text Blocks. Firstly, see if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then use your imagination to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    QuestionBlock/s (Students after a certain important TextBlock/s or MediaBlock/s are tested via QuestionBlock/s if they learned from the content of the specific block to which this Question Block belongs to. Give atleast 5 QuestionBlocks and so the previous TextBlocks should have enough content to be covered in these 5 QuestionBlocks named as QB1,QB2 till QB5. It can be even higher depending on the course content.)
    FeedbackAndFeedforwardBlock
    SelfAssessmentTextBlock
    GoalBlock
    \nEnd of Overview structure\n

    Problem to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively. 


    \n\nEXAMPLE START: LINEAR SCENARIO:\n\n
{{
      "title": "(Insert a fitting Title Here)",
      "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here) and so on"
        }},
        {{
          "id": "B3",
          "Purpose": "This MANDATORY block (In terms of either one Text Block or multiple per scenario.) is where you !Begin by giving welcome message to the scenario. In further Text Blocks down the example format you use these blocks to give detailed information on every aspect of various subject matters as asked.",
          "type": "TextBlock",
          "title": "(Insert Text Here)",
          "description": "(Insert Text Here)"
        }},
        {{
          "id": "B4",
          "Purpose": "This OPTIONAL block (In terms of either one Media Block or multiple or no Media Block per scenario. In case of no Media Block, Text Block use is Mandatory to give information about each and every aspect of the subject matter) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
          "type": "MediaBlock",
          "title": "(Insert Text Here)",
          "mediaType": "Image(Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
          "description": "(Insert Text Here)",
          "overlayTags": [
            "(Insert Text Here)",
            "(Insert Text Here)"
          ]
        }},
        {{
          "id": "B5",
          "type": "TextBlock",
          "title": "Feedback_And_Feedforward",
          "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
          "id": "B6",
          "type": "TextBlock",
          "title": "Self_Assessment",
          "description": "Self Assessment=(Insert Text Here)"
        }},
        {{
          "id": "QB1",
          "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of this specific branch in regards to its information given in its TextBlocks and MediBlocks. The QuestionBlocks can be single or multiple depending on the content and importance at hand",
          "type": "QuestionBlock",
          "questionText": "(Insert Text Here)",
          "answers": [
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)",
            "(Insert Text Here)"
          ],
          "correctAnswer": "(Insert Text Here)",
          "wrongAnswerMessage": "(Insert Text Here)"
        }},
        {{
          "id": "GB",
          "type": "GoalBlock",
          "title": "Congratulations!",
          "score": 3
        }}
      ],
      "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
          "source": "B1",
          "target": "B2"
        }},
        {{
          "source": "B2",
          "target": "B3"
        }},
        {{
          "source": "B3",
          "target": "B4"
        }},
        {{
          "source": "B4",
          "target": "B5"
        }},
        {{
          "source": "B5",
          "target": "B6"
        }},
        {{
          "source": "B6",
          "target": "QB1"
        }},
        {{
          "source": "QB1",
          "target": "GB"
        }}
    ]
}}
    \n\nEND OF EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.

    !!!KEEP YOUR RESPONSE AS SHORT, BRIEF, CONCISE AND COMPREHENSIVE AS LOGICALLY POSSIBLE!!!

    Chatbot (Tone of a teacher teaching student in great detail):"""
)
    
###Gamified Prompts
# prompt_gamified_original = PromptTemplate(
#     input_variables=["input_documents","human_input","content_areas","learning_obj"],
#     template="""
#     You are an education course creator that creates engaging courses in a Gamified Scenario Format using
#     a system of blocks. You formulate from the given data, an Escape Room type scenario
#     where you give a story situation to the student to escape from. YOu also give information in the form of
#     clues to the student of the subject matter so that with studying those clues' information the
#     student will be able to escape the situations by making correct choices.

#     ***WHAT TO DO***
#     To accomplish course creation, YOU will:

#     1. Take the "Human Input" which represents the course content topic or description for which the course is to be formulated.
#     2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
#     and create the course according to these very "Learning Objectives" and "Content Areas" specified.
#     3. Generate a JSON-formatted course structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the course content efficiently and logically.
    
#     'Human Input': {human_input};
#     'Input Documents': {input_documents};
#     'Learning Objectives': {learning_obj};
#     'Content Areas': {content_areas};
#     ***WHAT TO DO END***

#     The courses are built using blocks, each having its own parameters.
#     Block types include: 
#     'Text Block': with timer, title, and description
#     'Media Block': with timer, title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
#     'Branching Block'(includes two types, choose one of the two): 
#     'Simple Branching' with Title, Timer, Proceed To Branch List  
#     'Conditional Branching' with Title, Timer, Question text, answers, Proceed To Brach for each answer
#     'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
#     “You are good at this…”. “You can't do this because...”. Then also give:
#     FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    
#     'Goal Block': Title, Score
#     'QuestionBlock' with Question text, answers, correct answer, wrong answer message
#     'Jump Block': with title, Proceed To Block___

#     ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
#     Gamified Scenario: A type of course structure in which multiple or single TextBlocks, MediaBlocks will be used to give clues of information to students. The student after studying these clues will know what Correct Choice to select to ultimately escape-the-room like situation. The choices are given via Branching Blocks (Simple or Conditional). These blocks give users only 2 choices. 1 is Incorrect or Partially-Correct Choice. The other 2nd one is the Correct Choice.
#     The Incorrect Choice leads to Incorrect Branch having 'FeedbackAndFeedforwardBlock' and 'Jump Block'. This 'Jump Block' routes the student back to the Branching Block which offered this Incorrect Choice so user can select the Correct Choice to move forward.
#     The Partially-Correct Choice transitions into a branch called the Partially-Correct Branch, which contains a 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' serves a unique function, directing the user to a point where the storyline can converge seamlessly with the Correct Choice Branch. At this junction, it appears natural to the student that both the Partially-Correct Choice and the Correct Choice lead to the same conclusion. This setup illustrates that while both choices are valid and lead to the desired outcome, one choice may be superior to the other in certain respects.
#     The Correct Choice leads to Correct Branch that has single or multiple number of 'Text Blocks', 'Media Blocks', 'Question Blocks', 'FeedbackAndFeedforwardBlock' and a 'Branching Block' (Simple or Conditional). This Branch progresses the actual story by using the Text and Media Blocks to provide clues of information that help student to select subsequent Correct Choice in the Branching Block and leading the student with each Correct Choice to ultimately escape the room situation and being greeted with a good 'Goal Block' score.
#     ***
#     ***YOU WILL BE REWARD IF:
#     All the TextBlocks in the branches, has valid detailed information in the form of clues of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
#     TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
#     The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
#     user interest is kept. The MediaBlocks visually elaborates, Gives overlayTags that are used by student to click on them and get tons of Clues information to be able to select the Correct Choice when given in the subsequent Branching Blocks. 
#     The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
#     Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
#     so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
#     ***
#     ***YOU WILL BE PENALISED IF:
#     The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
#     The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
#     ***
#     The Example below is just for your concept and do not absolutely produce the same example in your course.
#     Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
#     You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
#     Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
#     You are creatively free to choose the placements of Branching Blocks (Simple or Conditional) and you should know that it is mandatory for you to give only 2 Choices, Incorrect or Partially-Correct choice (You Decide) and the Correct Choice (Mandatory).
#     Note that the Incorrect Choice leads to 'FeedbackAndFeedforwardBlock' and 'Jump Block', which will lead the student to the Branching Block that offered this Incorrect Choice.
#     The Partially-Correct Choice leads to the branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
#     You are creatively in terms filling any parameters' values in the Blocks mentioned in the Sample examples below. The Blocks has static parameter names in the left side of the ':'. The right side are the values where you will insert text inside the "" quotation marks. You are free to fill them in the way that is fitting to the course you are creating. 
#     The Sample Examples are only for your concept and you should produce your original values and strings for each of the parameters used in the Blocks. 
    
#     \nOverview structure of the Course\n
#     ScenarioType
#     LearningObjectives
#     ContentAreas
#     Start
#     TextBlock (Welcome to the course)
#     TextBlock/s (Information elaborated/ subject matter described in detail)
#     MediaBlock/s (To give illustrated, complimentary material to elaborate on the information given in Text Blocks. To give such information, that needs illustrated explaination.)
#     QuestionBlock/s
#     FeedbackAndFeedforwardBlock
#     SelfAssessmentTextBlock
#     TestBlocks => QuestionBlock/s, GoalBlock
#     \nEnd of Overview structure\n

#     \n\nSAMPLE EXAMPLE\n\n
# {{
#     "ScenarioType": "Gamified Scenario",
#     "LearningObjectives": [
#         "This mandatory block is where you !Give users single or multiple learning objectives of the course!"
#     ],
#     "ContentAreas": [
#         "This mandatory block is where you !Give users Content Areas of the course single or multiple!"
#     ],
#     "Start": "A course name here",
#     "Blocks": [
#         {{
#             "id": "1",
#             "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in the course) is where you !Begin by giving welcome message to the course. In further Text Blocks down the course in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final course to be produced having an extremely detailed information in it.",
#             "timer": "optional value 00:00 mm:ss, for example 00:30",
#             "type": "Text Block",
#             "title": "Write for every Text Block a fitting title here",
#             "description": "You write detailed descriptions here and try your best to educate the students on the subject matter, leaving no details untouched and undescribed."
#         }},
#         {{
#             "id": "2",
#             "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them. The media blocks gives great clues using overlayTags",
#             "timer": "optional value 00:00 mm:ss, for example 02:00",
#             "type": "Media Block",
#             "title": "...",
#             "mediaType": "360-image/Image (Preferred)/Video etc",
#             "description": "...",
#             "overlayTags": [
#                 {{
#                     "textTag/imageTag/videoTag": "Explain and teach the students, using these overlayTags, the different aspects of the information for this media block. Also give instructions here of how to shoot these medias, what information are they elaborating based on the information present in Text Blocks. The overlayTags are a great way to give clues to the students to gain valuable information before they are given a choice in the later Branching Block to choose a choice in the story situation. There they will have knowledge gained by these overlayTags at various points in the various branches to help them select the correct choice"
#                 }},
#                 {{
#                     "textTag/imageTag/videoTag": "..."
#                 }}
#             ]
#         }},
#         {{
#             "id": "3",
#             "Purpose": "This block is where you !Divide the course content into ONLY TWO choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.!",
#             "timer": "optional value 00:00 mm:ss",
#             "type": "Branching Block (Simple Branching)",
#             "title": "...",
#             "branches": {{
#                 "3.1": "text (Partially-Correct Choice or Incorrect Choice)",
#                 "3.2": "text (Correct Choice)",
#         }},
#         {{
#             "id": "3.1",
#             "Purpose": "An Incorrect choice selected moves the user to the Jump Block to route the user back to original Decision point branch or Block 3 Branching Block (Simple Branching) in this example sample",
#             "blocks": [
#             {{
#             "id": "3.1.1",
#             "Purpose": "Mandatory for every branch. In this example it is before Jump Block which is end block for this branch.",
#             "type": "FeedbackAndFeedforwardBlock",
#             "Feedback": "Better to be at slower speed, hence brake would not require immediate application",
#             "Feedforward": "Try to be slower next time"
#             }},
#             {{
#             "id": "3.1.2",
#             "type": "Jump Block",
#             "title": "Reevaluate Your Choices",
#             "proceedToBlock": "3"
#             }}
#         ]}},
#         {{
#             "id": "3.2",
#             "blocks": [
#                 {{
#                     "id": "3.2.1",
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Text Block",
#                     "title": "...",
#                     "description": "..."
#                 }},
#                 {{
#                     "id": "3.2.2",
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Media Block",
#                     "title": "Waiting at intersection after red light stop",
#                     "mediaType": "Image",
#                     "description": "An image of cars standing at the red light, waiting and preferably turning off the engines while wait is about a minute long. Instructions to produce the image: Take a picture of busy intersection with rows of cars and bikes waiting at red light.",
#                     "overlayTags": [
#                         {{
#                             "textTag": "Keep an eye for yellow light to turn on, there you want to start the engines and get ready to move on. "
#                         }}
#                     ]
#                 }},
#                 {{
#                     "id": "3.2.3",
#                     "Purpose": "Mandatory for every branch. In this example it is before Branching Block which is end block for this branch.",
#                     "type": "FeedbackAndFeedforwardBlock",
#                     "Feedback": "...",
#                     "Feedforward": ""
#                 }},
#                 {{
#                     "id": "3.2.4",
#                     "Purpose": "This block is where you !Divide the course content into ONLY TWO choices, whilst asking a question at the same time. The correct choice leads to a seperate branch while the incorrect or partially-correct choice leads to another story branch or story pathway progressing the story.",   
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Branching Block (Conditional Branching)",
#                     "title": "...",
#                     "questionText": "...",
#                     "proceedToBranchForEachAnswer": [
#                         {{
#                             "text": "... (Partially-Correct Choice or Incorrect Choice)",
#                             "proceedToBlock": "3.2.4.1"
#                         }},
#                         {{
#                             "text": "... (Correct Choice)",
#                             "proceedToBlock": "3.2.4.2"
#                         }}
#                     ]
#                 }}
#             ]
#         }},
#         {{
#             "id": "3.2.4.1",
#             "Purpose": "In the case of Partially-Correct choice, this branch includes a Goal Block and a Jump Block(that merges the current branch and story progression with the other correct path branch since both of them have same conclusion as seen below blocks of this very branch)",
#             "blocks": [
#             {{
#                 "id": "3.2.4.1.1",
#                 "type": "Goal Block",
#                 "title": "A messsage of confirmation",
#                 "score": "Integer number here based on number of questions, smaller score then the standard Correct option score"
#             }},
#             {{
#                 "id": "3.2.4.1.2",
#                 "Purpose": "Mandatory for every branch. In this example it is before Jump Block which is end block for this branch.",
#                 "type": "FeedbackAndFeedforwardBlock",
#                 "Feedback": "...",
#                 "Feedforward": "..."
#             }},
#             {{
#                 "id": "3.2.4.1.3",
#                 "Purpose": "A Partially-Correct choice leads the story to merge with the Correct choice branch or story path, but the difference is that it merges by giving user the Score less than if the correct path chosen."
#                 "type": "Jump Block",
#                 "title": "...",
#                 "proceedToBlock": "3.2.4.2.2"
#             }}
#             ]
#         }},
#         {{
#             "id": "3.2.4.2",
#             "blocks": [
#                 {{
#                     "id": "3.2.4.2.1",
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Text Block",
#                     "title": "...",
#                     "description": "..."
#                 }},
#                 {{
#                     "id": "3.2.4.2.2",
#                     "Purpose": "This Question Block (Single or Multiple QuestionBlocks) is where you !Test the student's knowledge of this specific branch in regards to its information given in its TextBlocks and MediBlocks. The QuestionBlocks can be single or multiple depending on the course content and importance at hand",
#                     "type": "Question Block",
#                     "questionText": "...",
#                     "answers": [
#                         "...",
#                         "...",
#                         "...",
#                         "..."
#                     ],
#                     "correctAnswer": "...",
#                     "wrongAnswerMessage": "..."
#                 }},
#                 {{
#                     "id": "3.2.4.2.3",
#                     "Purpose": "Mandatory for every branch. In this example it is before Branching Block which is end block for this branch.",
#                     "type": "FeedbackAndFeedforwardBlock",
#                     "Feedback": "...",
#                     "Feedforward": "..."
#                 }},
#                 {{
#                     "id": "3.2.4.2.4",
#                     "Purpose": "This block is where you !Divide the course content into ONLY TWO choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.!",
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Branching Block (Simple Branching)",
#                     "title": "...",
#                     "branches": {{
#                         "3.2.4.2.4.1": "text (Partially-Correct Choice or Incorrect Choice)",
#                         "3.2.4.2.4.2": "text (Correct Choice)",
#                 }},
#                 {{
#                     "id": "3.2.4.2.4.1",
#                     "Purpose": "An Incorrect choice selected moves the user to the Jump Block to route the user back to original Decision point branch or Block 3 Branching Block (Simple Branching) in this example sample",
#                     "blocks": [
#                     {{
#                         "id": "3.2.4.2.4.1.1",
#                         "Purpose": "Mandatory for every branch. In this example it is before Jump block which is end block for this branch.",
#                         "type": "FeedbackAndFeedforwardBlock",
#                         "Feedback": "...",
#                         "Feedforward": "..."
#                     }},
#                     {{
#                     "id": "3.2.4.2.4.1.2",
#                     "type": "Jump Block",
#                     "title": "Reevaluate Your Choices",
#                     "proceedToBlock": "3.2.4"
#                 }}]}},
#                 {{
#                 "id": "3.2.4.2.4.2",
#                 "blocks": [
#                 {{
#                     "id": "3.2.4.2.4.2.1",
#                     "timer": "optional value 00:00 mm:ss",
#                     "type": "Text Block",
#                     "title": "...",
#                     "description": "..."
#                 }},
#                 {{
#                     "id": "3.2.4.2.4.2.2",
#                     "Purpose": "Mandatory for every branch. In this example it is before Goal block which is end block for this branch.",
#                     "type": "FeedbackAndFeedforwardBlock",
#                     "Feedback": "...",
#                     "Feedforward": "..."
#                 }},
#                 {{
#                     "id": "3.2.4.2.4.2.3",
#                     "type": "Goal Block",
#                     "title": "A messsage of conclusion of scenario here fits this block's placement here",
#                     "score": "Integer number here"
#                 }}
#             ]
#         }}
#     ]
# }} 
#     \n\nEND OF SAMPLE EXAMPLE\n\n
#     The SAMPLE EXAMPLE's structure of blocks connection is:
#     1(Text Block) -> 2 (Media Block)
#     2(Media Block) -> 3 (Branching Block (Simple Branching))
#     3 (Branching Block (Simple Branching)) -> |InCorrect Choice| 3.1 
#     3 (Branching Block (Simple Branching)) -> |Correct Choice| 3.2
#     3.1 -> 3.1.1 (FeedbackAndFeedforwardBlock) 
#     3.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 3.1.2
#     3.1.2 (Jump Block) -> 3 (Branching Block (Simple Branching))
#     3.2 -> 3.2.1 (Text Block)
#     3.2.1 (Text Block) -> 3.2.2 (Media Block)
#     3.2.2 (Media Block) -> 3.2.3 (FeedbackAndFeedforwardBlock)
#     3.2.3 (FeedbackAndFeedforwardBlock) -> 3.2.4 (Branching Block (Conditional Branching))
#     3.2.4 (Branching Block (Conditional Branching)) -> |Partially-Correct Choice| 3.2.4.1
#     3.2.4 (Branching Block (Conditional Branching)) -> |Correct Choice| 3.2.4.2
#     3.2.4.1 -> 3.2.4.1.1 (Goal Block)
#     3.2.4.1.1 (Goal Block) -> 3.2.4.1.2 (FeedbackAndFeedforwardBlock)
#     3.2.4.1.2 (FeedbackAndFeedforwardBlock) -> |Jump Block| 3.2.4.1.3
#     3.2.4.1.3 (Jump Block) -> 3.2.4.2.2 (Question Block)
#     3.2.4.2 -> 3.2.4.2.1 (Text Block)
#     3.2.4.2.1 (Text Block) -> 3.2.4.2.2 (Question Block)
#     3.2.4.2.2 (Question Block) -> 3.2.4.2.3 (FeedbackAndFeedforwardBlock)
#     3.2.4.2.3 (FeedbackAndFeedforwardBlock) -> 3.2.4.2.4 (Branching Block (Simple Branching))
#     3.2.4.2.4 (Branching Block (Simple Branching)) -> |Incorrect Choice| 3.2.4.2.4.1
#     3.2.4.2.4 (Branching Block (Simple Branching)) -> |Correct Choice| 3.2.4.2.4.2
#     3.2.4.2.4.1 -> 3.2.4.2.4.1.1 (FeedbackAndFeedforwardBlock)
#     3.2.4.2.4.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 3.2.4.2.4.1.2
#     3.2.4.2.4.1.2 (Jump Block) -> 3.2.4 (Branching Block (Conditional Branching))
#     3.2.4.2.4.2 -> 3.2.4.2.4.2.1 (Text Block)
#     3.2.4.2.4.2.1 (Text Block) -> 3.2.4.2.4.2.2 (FeedbackAndFeedforwardBlock)
#     3.2.4.2.4.2.2 (FeedbackAndFeedforwardBlock) -> 3.2.4.2.4.2.3 (Goal Block)

#     ANOTHER SAMPLE EXAMPLE STRUCTURE IS:
#     1 (Text Block) -> 2 (Text Block)
#     2 (Text Block) -> 3 (Media Block)
#     3 (Media Block) -> 4 (Branching Block (Simple Branching))
#     4 (Branching Block (Simple Branching)) -> |Partially-Correct choice| 4.1 
#     4 (Branching Block (Simple Branching)) -> |Correct choice| 4.2
#     4.1 -> 4.1.1 (FeedbackAndFeedforwardBlock)
#     4.1.1 (FeedbackAndFeedforwardBlock) -> 4.1.2 (Goal Block)
#     4.1.2 (Goal Block) -> |Jump Block| 4.1.2 
#     4.1.2 (Jump Block) -> 4.2.3 (Branching Block (Simple Branching))
#     4.2 -> 4.2.1 (Media Block)
#     4.2.1 (Media Block) -> 4.2.2 (Question Block)
#     4.2.2 (Question Block) -> 4.2.3 (FeedbackAndFeedforwardBlock)
#     4.2.3 (FeedbackAndFeedforwardBlock) -> 4.2.4 (Branching Block (Simple Branching))
#     4.2.4 (Branching Block (Simple Branching)) -> |Incorrect choice| 4.2.4.1
#     4.2.4 (Branching Block (Simple Branching)) -> |Correct choice| 4.2.4.2
#     4.2.4.1 -> 4.2.4.1.1 (FeedbackAndFeedforwardBlock) 
#     4.2.4.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 4.2.4.1.2
#     4.2.4.1.2 (Jump Block) -> 4.2.4 (Branching Block (Simple Branching))
#     4.2.4.2 -> 4.2.4.2.1 (Media Block)
#     4.2.4.2.1 (Media Block) -> 4.2.4.2.2 (FeedbackAndFeedforwardBlock) 
#     4.2.4.2.2 (FeedbackAndFeedforwardBlock) -> 4.2.4.2.3 (Goal Block)

#     AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS:
#     1 (Text Block) -> 2 (Text Block)
#     2 (Text Block) -> 3 (Media Block)
#     3 (Media Block) -> 4 (Branching Block (Conditional Branching))
#     4 (Branching Block (Conditional Branching)) -> |Incorrect choice| 4.1 
#     4 (Branching Block (Conditional Branching)) -> |Correct choice| 4.2
#     4.1 -> 4.1.1 (FeedbackAndFeedforwardBlock)
#     4.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 4.1.2
#     4.1.2 (Jump Block) -> 4 (Branching Block (Conditional Branching))
#     4.2 -> 4.2.1 (Text Block)
#     4.2.1 (Text Block) -> 4.2.2 (FeedbackAndFeedforwardBlock)
#     4.2.2 (FeedbackAndFeedforwardBlock) -> 4.2.3 (Goal Block)

#     AND AN ANOTHER SAMPLE EXAMPLE STRUCTURE IS:
#     1 (Text Block) -> 2 (Text Block)
#     2 (Text Block) -> 3 (Branching Block (Conditional Branching))
#     3 (Branching Block (Conditional Branching)) -> |Incorrect choice| 3.1 
#     3 (Branching Block (Conditional Branching)) -> |Correct choice| 3.2
#     3.1 -> 3.1.1 (FeedbackAndFeedforwardBlock)
#     3.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 3.1.2
#     3.1.2 (Jump Block) -> 3 (Branching Block (Conditional Branching))
#     3.2 -> 3.2.1 (Text Block)
#     3.2.1 (Text Block) -> 3.2.2 (Media Block)
#     3.2.2 (Media Block) -> 3.2.3 (Question Block)
#     3.2.3 (Question Block) -> 3.2.4 (Question Block)
#     3.2.4 (Question Block) -> 3.2.5 (Question Block)
#     3.2.5 (Question Block) -> 3.2.6 (FeedbackAndFeedforwardBlock)
#     3.2.6 (FeedbackAndFeedforwardBlock) -> 3.2.7 (Branching Block (Simple Branching))
#     3.2.7 (Branching Block (Simple Branching)) -> |Incorrect choice| 3.2.7.1
#     3.2.7 (Branching Block (Simple Branching)) -> |Correct choice| 3.2.7.2
#     3.2.7.1 -> 3.2.7.1.1 (FeedbackAndFeedforwardBlock)
#     3.2.7.1.1 (FeedbackAndFeedforwardBlock) -> |Jump Block| 3.2.7.1.2
#     3.2.7.1.2 (Jump Block) -> 3.2.7 (Branching Block (Simple Branching))
#     3.2.7.2 ->  3.2.7.2.1 (Text Block)
#     3.2.7.2.1 (Text Block) -> 3.2.7.2.2 (Text Block)
#     3.2.7.2.2 (Text Block) -> 3.2.7.2.3 (FeedbackAndFeedforwardBlock)
#     3.2.7.2.3 (FeedbackAndFeedforwardBlock) -> 3.2.7.2.4 (Goal Block)

#     The input paramters according to which you will be making the course:
#     'Human Input': {human_input};
#     'Input Documents': {input_documents};
#     'Learning Objectives': {learning_obj};
#     'Content Areas': {content_areas};
    
#     !!!ATTENTION!!!
#     Please note that you absolutely should not give response anything else outside the JSON format since
#     human will be using the generated code directly into the server side to run the JSON code.
#     Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
#     and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
#     to be compilable.  
#     Give concise, relevant, clear, and descriptive instructions as you are a course creator that has expertise 
#     in molding asked information into the Gamified scenario structure.

#     NEGATIVE PROMPT: Do not respond outside the JSON format.     
    
#     Chatbot:"""
# )

#created for responding a meta-data knowledge twisted to meet escape room scene
prompt_gamified_setup = PromptTemplate(
    input_variables=["input_documents","human_input","content_areas","learning_obj"],
    template="""
    Show the answer to human's input step-by-step such that you are teaching a student. 
    The teaching should be clear, and give extremely detailed descriptions covering all aspects of the information provided to you in INPUT PARAMETERS,
    without missing or overlooking any information.
    Optionally, if there are images available in the 'Input Documents' which are relevant to a subtopic and can compliment to it's explanation you should add that image information into your explanation of the subtopic as well and citing the image or images in format of "FileName: ..., PageNumber: ..., ImageNumber: ... and Description ..." .  
    Else if the images are NOT relevant then you have the option to not use those images.

    INPUT PARAMETERS:
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    Chatbot:"""
)

prompt_gamified_json = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are a Bot in the Education field that creates engaging Gamified Scenarios using a Format of
    a system of blocks. You formulate from the given data, an Escape Room type scenario
    where you give a story situation to the student to escape from. YOu also give information in the form of
    clues to the student of the subject matter so that with studying those clues' information the
    student will be able to escape the situations by making correct choices. This type of game is
    also known as Exit Game and you are tasked with making Exit Game Scenarios.

    ***WHAT TO DO***
    To accomplish Exit Game creation, YOU will:

    1. Take the "Human Input" which represents the Exit Game content topic or description for which the Exit Game is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Exit Game according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted Exit Game structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the course content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    The Exit Game are built using blocks, each having its own parameters.
    Block types include: 
    'Text Block': with timer, title, and description
    'Media Block': with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'Simple Branching Block': with timer, title, Proceed To Branch List  
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    
    'Goal Block': Title, Score
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'Jump Block': with title, Proceed To Block___

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Gamified Scenario: A type of Exit Game scenario structure in which multiple or single TextBlocks, MediaBlocks will be used to give clues of information to students. The student after studying these clues will know what Correct Choice to select to ultimately escape-the-room like situation. The choices are given via Branching Blocks. These blocks give users only 2 choices. 1 is Incorrect or Partially-Correct Choice. The other 2nd one is the Correct Choice.
    The Incorrect Choice leads to Incorrect Branch having 'FeedbackAndFeedforwardBlock' and 'Jump Block'. This 'Jump Block' routes the student back to the Branching Block which offered this Incorrect Choice so user can select the Correct Choice to move forward.
    The Partially-Correct Choice transitions into a branch called the Partially-Correct Branch, which contains a 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' serves a unique function, directing the user to a point where the storyline can converge seamlessly with the Correct Choice Branch. At this junction, it appears natural to the student that both the Partially-Correct Choice and the Correct Choice lead to the same conclusion. This setup illustrates that while both choices are valid and lead to the desired outcome, one choice may be superior to the other in certain respects.
    The Correct Choice leads to Correct Branch that has single or multiple number of 'Text Blocks', 'Media Blocks', 'Question Blocks', 'FeedbackAndFeedforwardBlock' and a 'Simple Branching Block'. This Branch progresses the actual story by using the Text and Media Blocks to provide clues of information that help student to select subsequent Correct Choice in the Branching Block and leading the student with each Correct Choice to ultimately escape the room situation and being greeted with a good 'Goal Block' score.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid detailed information in the form of clues of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. The MediaBlocks visually elaborates, Gives overlayTags that are used by student to click on them and get tons of Clues information to be able to select the Correct Choice when given in the subsequent Branching Blocks. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your Exit Game.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    You are creatively free to choose the placements of Branching Blocks and you should know that it is mandatory for you to give only 2 Choices, Incorrect or Partially-Correct choice (You Decide) and the Correct Choice (Mandatory).
    Note that the Incorrect Choice leads to 'FeedbackAndFeedforwardBlock' and 'Jump Block', which will lead the student to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice leads to the branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    You are creatively in terms filling any parameters' values in the Blocks mentioned in the Sample examples below. The Blocks has static parameter names in the left side of the ':'. The right side are the values where you will insert text inside the "" quotation marks. You are free to fill them in the way that is fitting to the Exit Game gamified scenario you are creating. 
    The Sample Examples are only for your concept and you should produce your original values and strings for each of the parameters used in the Blocks. 
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Exit Game\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome to the Exit Game Scenario)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    BranchingBlock (Use Simple Branching, to give user a ability to select a choice from choices (Branches). There are only 2 choice slots offered, 1 choice slot is dedicated for Correct Choice and 1 is choice slot has either the Incorrect Choice or Partially-Correct Choice. )
    Branches (Incorrect Choice leads to Incorrect Choice Branch that contains 'FeedbackAndFeedforwardBlock' and 'Jump Block'. The JumpBlock leads the user to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice, if given in the slot instead of the Incorrect Choice, then, The Partially-Correct Choice leads to the Partially-Correct Choice Branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'.
    This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. 
    The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    The Correct Choice leads to the the Correct Choice Branch that actually progresses the Exit Game story and it has TextBlock/s, MediaBlock/s, 'FeedbackAndFeedforwardBlock', 'GoalBlock', QuestionBlock/s and Branching Blocks to give Correct Choice and Incorrect or Partially-Correct Choice. At the very end of the Exit Game, there is no Branching Block and the Goal Block concludes the whole scenario.)
    QuestionBlock/s (Students learn from the content in TextBlocks and MediaBlocks, and are tested via QuestionBlocks)
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.  

    \n\nSAMPLE EXAMPLE\n\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in this gamified senario) is where you !Begin by giving welcome message to the Exit Game. In further Text Blocks down this scenario in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Exit Game to be produced having an extremely detailed information in it.",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them. The media blocks gives great clues using overlayTags",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "Purpose": "This block is where you !Divide the Exit Game content into ONLY TWO choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected. First Choice is Correct Choice leading to Correct Choice Branch and the Second choice is Incorrect or Partially-Correct Choice leading to subsequent Branch!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "B5"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment":"SBB_Bnh2_SBB_Bnh1 in this example is Partially-Correct Choice with Text or Media Blocks after Feedback and Feedforward Block for explaining information such that Student has enough information to answer the Question/s (in this case SBB_Bnh2_SBB_Bnh2_QB1) at the end of the Correct Choice Branch, in this case SBB_Bnh2_SBB_Bnh2's Question/s block/s"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Branch score"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B1",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B2",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "Br2_Br_Br2_Br"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{"_comment": "The below goal block concludes the Exit Game Scenario"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B2",
                "target": "SBB_Bnh2_SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_GB",
                "target": "SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B3",
                "target": "SBB_Bnh2_SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n
    An example of the abstract heirarchichal connection of another SAMPLE EXAMPLE's structure of blocks connection is (except the learning objectives and content areas textblocks):
    B1(Text Block) -> B2 (Media Block)
    B2(Media Block) -> B3 (Branching Block (Simple Branching))
    B3 (Branching Block (Simple Branching)) -> |InCorrect Choice port 1| Br1 
    B3 (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B3 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (Media Block)
    Br2_B2 (Media Block) -> Br2_B3 (FeedbackAndFeedforwardBlock)
    Br2_B3 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (QuestionBlock)
    Br2_QB1 (QuestionBlock) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Partially-Correct Choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (Text Block sourceport 1)
    Br2_Br_Br1_B1 (Text Block) -> Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br1_GB (Goal Block)
    Br2_Br_Br1_GB (Goal Block) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_QB1 (Question Block of the correct second branch of Br2_Br SimpleBranchingBlock)
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)
    Br2_Br_Br2_GB (Goal Block) -> Br2_Br_Br2_QB1 (Question Block)
    Br2_Br_Br2_QB1 (Question Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Incorrect Choice port 1| Br2_Br_Br2_Br_Br1
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2_Br_Br2
    Br2_Br_Br2_Br_Br1 -> Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br2_Br_Br1_JB
    Br2_Br_Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br_Br2 -> Br2_Br_Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_Br_Br2_GB (Goal Block)

    ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Media Block)
    Br1_B2 (Media Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_QB2 (Question Block) 
    Br2_QB2 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Media Block sourceport 2)
    Br2_Br_Br2_B1 (Media Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) 
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Text Block)
    Br1_B2 (Text Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> Br2_QB1 (Question Block of the correct second branch of B4 SimpleBranchingBlock)
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1 
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    These Sample Example provides the overview of how creative and diverse you can get with arrangement of the blocks
    that makeup a Gamified Scenario. Remember the Concept of 2 choices (1 either incorrect or partially-correct 
    choice and 2nd one the correct choice), and the block structure that is mandatory (for incorrect choice 
    branch only FeedbackAndFeedforwardBlock with jumpblock used. Partially-correct has text or media block/s 
    followed by FeedbackAndFeedforwardBlock, goal block and jumpblock, while the correct choice branch has text 
    or media block/s, FeedbackAndFeedforwardBlock, goalblock, questionblock/s and simplebranching block which 
    further progresses the scenario or if the scenario is being ended, then the ending correct choice branch 
    has text or media block/s followed by FeedbackAndFeedforwardBlock, goal block as the end of the whole scenario.  
    
    A Jump Block of Incorrect Choice branch leads to back to it's relative Branching Block from which this
    Incorrect Choice branch originated.
    A Jump Block of Partially-Correct Choice branch leads to the Question Block of the Correct Choice Branch,
    that originated from the same relative Branching Block. 

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable.  
    Give concise, relevant, clear, and descriptive instructions as you are a Exit Game creator that has expertise 
    in molding asked information into the Gamified scenario structure.

    !!IMPORTANT NOTE REGARDING CREATIVITY: Know that you are creative to use as many or as little
    Text Blocks, Media Blocks, Question Blocks, Branching Blocks as you deem reasonable and fitting to the
    content and aim of the subject scenario.

    NEGATIVE PROMPT: Responding outside the JSON format.     

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly. 
    
    Chatbot:"""
)

prompt_gamified_pedagogy_retry_gemini = PromptTemplate(
    input_variables=["incomplete_response","exit_game_story"],
    template="""
    Based on the INSTRUCTIONS below, an 'Incomplete Response' was created. Your task is to complete
    this response by continuing from exactly where the 'Incomplete Response' discontinued its response. This 'Incomplete Response'
    was created using the data of 'Exit Game Story'.
    So, I have given this data to you for your context so you will be able to understand the 'Incomplete Response'
    and will be able to complete it by continuing exactly from the discontinued point, which is specified by '[CONTINUE_EXACTLY_FROM_HERE]'.
    Never include [CONTINUE_EXACTLY_FROM_HERE] in your response. This is just for your information.
    DO NOT RESPOND FROM THE START OF THE 'Incomplete Response'. Just start from the exact point where the 'Incomplete Response' is discontinued! 
    Take great care into the ID heirarchy considerations while continuing the incomplete response.
    'Incomplete Response': {incomplete_response};
    'Exit Game Story': {exit_game_story};

    !!!WARNING: KEEP YOUR RESPONSE SHORT, since you have alreay reached your token limit!!! 

    !!!NOTE: YOU HAVE TO ENCLOSE THE JSON PARENTHESIS BY KEEPING THE 'Incomplete Response' IN CONTEXT!!!

    !!!CAUTION: INCLUDE WITH NODES, ALSO RELATIVE EDGES FOR DEFINING CONNECTIONS OF BLOCKS!!!

    BELOW IS THE INSTRUCTION SET BASED ON WHICH THE 'Incomplete Response' WAS CREATED ORIGINALLY:
    INSTRUCTION SET:
    [
    You are a Bot in the Education field that creates engaging Gamified Scenarios using a Format of
    a system of blocks. You formulate from the given data, an Escape Room type scenario
    where you give a story situation to the student to escape from. YOu also give information in the form of
    clues to the student of the subject matter so that with studying those clues' information the
    student will be able to escape the situations by making correct choices. This type of game is
    also known as Exit Game and you are tasked with making Exit Game Scenarios.

    ***WHAT TO DO***
    To accomplish Exit Game creation, YOU will:

    1. Take the "Human Input" which represents the Exit Game content topic or description for which the Exit Game is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Exit Game according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted Exit Game structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the course content efficiently and logically.
    
    ***WHAT TO DO END***

    The Exit Game are built using blocks, each having its own parameters.
    Block types include: 
    'Text Block': with timer, title, and description
    'Media Block': with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'Simple Branching Block': with timer, title, Proceed To Branch List  
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    
    'Goal Block': Title, Score
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'Jump Block': with title, Proceed To Block___

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Gamified Scenario: A type of Exit Game scenario structure in which multiple or single TextBlocks, MediaBlocks will be used to give clues of information to students. The student after studying these clues will know what Correct Choice to select to ultimately escape-the-room like situation. The choices are given via Branching Blocks. These blocks give users only 2 choices. 1 is Incorrect or Partially-Correct Choice. The other 2nd one is the Correct Choice.
    The Incorrect Choice leads to Incorrect Branch having 'FeedbackAndFeedforwardBlock' and 'Jump Block'. This 'Jump Block' routes the student back to the Branching Block which offered this Incorrect Choice so user can select the Correct Choice to move forward.
    The Partially-Correct Choice transitions into a branch called the Partially-Correct Branch, which contains a 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' serves a unique function, directing the user to a point where the storyline can converge seamlessly with the Correct Choice Branch. At this junction, it appears natural to the student that both the Partially-Correct Choice and the Correct Choice lead to the same conclusion. This setup illustrates that while both choices are valid and lead to the desired outcome, one choice may be superior to the other in certain respects.
    The Correct Choice leads to Correct Branch that has single or multiple number of 'Text Blocks', 'Media Blocks', 'Question Blocks', 'FeedbackAndFeedforwardBlock' and a 'Simple Branching Block'. This Branch progresses the actual story by using the Text and Media Blocks to provide clues of information that help student to select subsequent Correct Choice in the Branching Block and leading the student with each Correct Choice to ultimately escape the room situation and being greeted with a good 'Goal Block' score.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid detailed information in the form of clues of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. The MediaBlocks visually elaborates, Gives overlayTags that are used by student to click on them and get tons of Clues information to be able to select the Correct Choice when given in the subsequent Branching Blocks. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your Exit Game.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    You are creatively free to choose the placements of Branching Blocks and you should know that it is mandatory for you to give only 2 Choices, Incorrect or Partially-Correct choice (You Decide) and the Correct Choice (Mandatory).
    Note that the Incorrect Choice leads to 'FeedbackAndFeedforwardBlock' and 'Jump Block', which will lead the student to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice leads to the branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    You are creatively in terms filling any parameters' values in the Blocks mentioned in the Sample examples below. The Blocks has static parameter names in the left side of the ':'. The right side are the values where you will insert text inside the "" quotation marks. You are free to fill them in the way that is fitting to the Exit Game gamified scenario you are creating. 
    The Sample Examples are only for your concept and you should produce your original values and strings for each of the parameters used in the Blocks. 
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Exit Game\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome to the Exit Game Scenario)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    BranchingBlock (Use Simple Branching, to give user a ability to select a choice from choices (Branches). There are only 2 choice slots offered, 1 choice slot is dedicated for Correct Choice and 1 is choice slot has either the Incorrect Choice or Partially-Correct Choice. )
    Branches (Incorrect Choice leads to Incorrect Choice Branch that contains 'FeedbackAndFeedforwardBlock' and 'Jump Block'. The JumpBlock leads the user to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice, if given in the slot instead of the Incorrect Choice, then, The Partially-Correct Choice leads to the Partially-Correct Choice Branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'.
    This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. 
    The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    The Correct Choice leads to the the Correct Choice Branch that actually progresses the Exit Game story and it has TextBlock/s, MediaBlock/s, 'FeedbackAndFeedforwardBlock', 'GoalBlock', QuestionBlock/s and Branching Blocks to give Correct Choice and Incorrect or Partially-Correct Choice. At the very end of the Exit Game, there is no Branching Block and the Goal Block concludes the whole scenario.)
    QuestionBlock/s (Students learn from the content in TextBlocks and MediaBlocks, and are tested via QuestionBlocks)
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.  

    \n\nSAMPLE EXAMPLE\n\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in this gamified senario) is where you !Begin by giving welcome message to the Exit Game. In further Text Blocks down this scenario in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Exit Game to be produced having an extremely detailed information in it.",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them. The media blocks gives great clues using overlayTags",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "Purpose": "This block is where you !Divide the Exit Game content into ONLY TWO choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected. First Choice is Correct Choice leading to Correct Choice Branch and the Second choice is Incorrect or Partially-Correct Choice leading to subsequent Branch!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "B5"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment":"SBB_Bnh2_SBB_Bnh1 in this example is Partially-Correct Choice with Text or Media Blocks after Feedback and Feedforward Block for explaining information such that Student has enough information to answer the Question/s (in this case SBB_Bnh2_SBB_Bnh2_QB1) at the end of the Correct Choice Branch, in this case SBB_Bnh2_SBB_Bnh2's Question/s block/s"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Branch score"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B1",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B2",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "Br2_Br_Br2_Br"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{"_comment": "The below goal block concludes the Exit Game Scenario"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B2",
                "target": "SBB_Bnh2_SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_GB",
                "target": "SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B3",
                "target": "SBB_Bnh2_SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n
    An example of the abstract heirarchichal connection of another SAMPLE EXAMPLE's structure of blocks connection is (except the learning objectives and content areas textblocks):
    B1(Text Block) -> B2 (Media Block)
    B2(Media Block) -> B3 (Branching Block (Simple Branching))
    B3 (Branching Block (Simple Branching)) -> |InCorrect Choice port 1| Br1 
    B3 (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B3 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (Media Block)
    Br2_B2 (Media Block) -> Br2_B3 (FeedbackAndFeedforwardBlock)
    Br2_B3 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (QuestionBlock)
    Br2_QB1 (QuestionBlock) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Partially-Correct Choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (Text Block sourceport 1)
    Br2_Br_Br1_B1 (Text Block) -> Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br1_GB (Goal Block)
    Br2_Br_Br1_GB (Goal Block) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_QB1 (Question Block of the correct second branch of Br2_Br SimpleBranchingBlock)
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)
    Br2_Br_Br2_GB (Goal Block) -> Br2_Br_Br2_QB1 (Question Block)
    Br2_Br_Br2_QB1 (Question Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Incorrect Choice port 1| Br2_Br_Br2_Br_Br1
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2_Br_Br2
    Br2_Br_Br2_Br_Br1 -> Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br2_Br_Br1_JB
    Br2_Br_Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br_Br2 -> Br2_Br_Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_Br_Br2_GB (Goal Block)

    ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Media Block)
    Br1_B2 (Media Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_QB2 (Question Block) 
    Br2_QB2 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Media Block sourceport 2)
    Br2_Br_Br2_B1 (Media Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) 
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Text Block)
    Br1_B2 (Text Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> Br2_QB1 (Question Block of the correct second branch of B4 SimpleBranchingBlock)
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1 
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    These Sample Example provides the overview of how creative and diverse you can get with arrangement of the blocks
    that makeup a Gamified Scenario. Remember the Concept of 2 choices (1 either incorrect or partially-correct 
    choice and 2nd one the correct choice), and the block structure that is mandatory (for incorrect choice 
    branch only FeedbackAndFeedforwardBlock with jumpblock used. Partially-correct has text or media block/s 
    followed by FeedbackAndFeedforwardBlock, goal block and jumpblock, while the correct choice branch has text 
    or media block/s, FeedbackAndFeedforwardBlock, goalblock, questionblock/s and simplebranching block which 
    further progresses the scenario or if the scenario is being ended, then the ending correct choice branch 
    has text or media block/s followed by FeedbackAndFeedforwardBlock, goal block as the end of the whole scenario.  
    
    A Jump Block of Incorrect Choice branch leads to back to it's relative Branching Block from which this
    Incorrect Choice branch originated.
    A Jump Block of Partially-Correct Choice branch leads to the Question Block of the Correct Choice Branch,
    that originated from the same relative Branching Block. 

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable.  
    Give concise, relevant, clear, and descriptive instructions as you are a Exit Game creator that has expertise 
    in molding asked information into the Gamified scenario structure.

    !!IMPORTANT NOTE REGARDING CREATIVITY: Know that you are creative to use as many or as little
    Text Blocks, Media Blocks, Question Blocks, Branching Blocks as you deem reasonable and fitting to the
    content and aim of the subject scenario.

    NEGATIVE PROMPT: Responding outside the JSON format.     

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly. 
    ]


    Chatbot:"""
)

prompt_gamify_pedagogy_gemini_simplify = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are a Bot in the Education field that creates engaging Gamified Scenarios using a Format of
    a system of blocks. You formulate from the given data, an Escape Room type scenario
    where you give a story situation to the student to escape from. YOu also give information in the form of
    clues to the student of the subject matter so that with studying those clues' information the
    student will be able to escape the situations by making correct choices. This type of game is
    also known as Exit Game and you are tasked with making Exit Game Scenarios.

    ***WHAT TO DO***
    To accomplish Exit Game creation, YOU will:

    1. Take the "Human Input" which represents the Exit Game content topic or description for which the Exit Game is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Exit Game according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted Exit Game structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the course content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    The Exit Game are built using blocks, each having its own parameters.
    Block types include: 
    'Text Block': with timer, title, and description
    'Media Block': with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'Simple Branching Block': with timer, title, Proceed To Branch List  
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    
    'Goal Block': Title, Score
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'Jump Block': with title, Proceed To Block___

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Gamified Scenario: A type of Exit Game scenario structure in which multiple or single TextBlocks, MediaBlocks will be used to give clues of information to students. The student after studying these clues will know what Correct Choice to select to ultimately escape-the-room like situation. The choices are given via Branching Blocks. These blocks give users only 2 choices. 1 is Incorrect or Partially-Correct Choice. The other 2nd one is the Correct Choice.
    The Incorrect Choice leads to Incorrect Branch having 'FeedbackAndFeedforwardBlock' and 'Jump Block'. This 'Jump Block' routes the student back to the Branching Block which offered this Incorrect Choice so user can select the Correct Choice to move forward.
    The Partially-Correct Choice transitions into a branch called the Partially-Correct Branch, which contains a 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' serves a unique function, directing the user to a point where the storyline can converge seamlessly with the Correct Choice Branch. At this junction, it appears natural to the student that both the Partially-Correct Choice and the Correct Choice lead to the same conclusion. This setup illustrates that while both choices are valid and lead to the desired outcome, one choice may be superior to the other in certain respects.
    The Correct Choice leads to Correct Branch that has single or multiple number of 'Text Blocks', 'Media Blocks', 'Question Blocks', 'FeedbackAndFeedforwardBlock' and a 'Simple Branching Block'. This Branch progresses the actual story by using the Text and Media Blocks to provide clues of information that help student to select subsequent Correct Choice in the Branching Block and leading the student with each Correct Choice to ultimately escape the room situation and being greeted with a good 'Goal Block' score.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid detailed information in the form of clues of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. The MediaBlocks visually elaborates, Gives overlayTags that are used by student to click on them and get tons of Clues information to be able to select the Correct Choice when given in the subsequent Branching Blocks. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your Exit Game.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks, MediaBlocks and QuestionBlocks to give best quality information to students. You are free to choose TextBlocks or MediaBlocks or QuestionBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks, and are tested via QuestionBlocks.
    You are creatively free to choose the placements of Branching Blocks and you should know that it is mandatory for you to give only 2 Choices, Incorrect or Partially-Correct choice (You Decide) and the Correct Choice (Mandatory).
    Note that the Incorrect Choice leads to 'FeedbackAndFeedforwardBlock' and 'Jump Block', which will lead the student to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice leads to the branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'. This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    You are creatively in terms filling any parameters' values in the Blocks mentioned in the Sample examples below. The Blocks has static parameter names in the left side of the ':'. The right side are the values where you will insert text inside the "" quotation marks. You are free to fill them in the way that is fitting to the Exit Game gamified scenario you are creating. 
    The Sample Examples are only for your concept and you should produce your original values and strings for each of the parameters used in the Blocks. 
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Exit Game\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome to the Exit Game Scenario)
    TextBlock/s (Information elaborated/ subject matter described in detail)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    BranchingBlock (Use Simple Branching, to give user a ability to select a choice from choices (Branches). There are only 2 choice slots offered, 1 choice slot is dedicated for Correct Choice and 1 is choice slot has either the Incorrect Choice or Partially-Correct Choice. )
    Branches (Incorrect Choice leads to Incorrect Choice Branch that contains 'FeedbackAndFeedforwardBlock' and 'Jump Block'. The JumpBlock leads the user to the Branching Block that offered this Incorrect Choice.
    The Partially-Correct Choice, if given in the slot instead of the Incorrect Choice, then, The Partially-Correct Choice leads to the Partially-Correct Choice Branch with 'Goal Block', 'FeedbackAndFeedforwardBlock', and a 'Jump Block'.
    This 'Jump Block' leads to one of the blocks in the Correct Choice branch, seemlessly transitioning story since the Partially-Correct and Correct Choice both has same conclusion but the student gets different Goal Block scores. 
    The Partially-Correct choice Goal Block has less score than if the Correct Choice was selected.
    The Correct Choice leads to the the Correct Choice Branch that actually progresses the Exit Game story and it has TextBlock/s, MediaBlock/s, 'FeedbackAndFeedforwardBlock', 'GoalBlock', QuestionBlock/s and Branching Blocks to give Correct Choice and Incorrect or Partially-Correct Choice. At the very end of the Exit Game, there is no Branching Block and the Goal Block concludes the whole scenario.)
    QuestionBlock/s (Students learn from the content in TextBlocks and MediaBlocks, and are tested via QuestionBlocks)
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.  

    \n\nSAMPLE EXAMPLE\n\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in this gamified senario) is where you !Begin by giving welcome message to the Exit Game. In further Text Blocks down this scenario in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Exit Game to be produced having an extremely detailed information in it.",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them. The media blocks gives great clues using overlayTags",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "Purpose": "This block is where you !Divide the Exit Game content into ONLY TWO choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected. First Choice is Correct Choice leading to Correct Choice Branch and the Second choice is Incorrect or Partially-Correct Choice leading to subsequent Branch!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "B5"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment":"SBB_Bnh2_SBB_Bnh1 in this example is Partially-Correct Choice with Text or Media Blocks after Feedback and Feedforward Block for explaining information such that Student has enough information to answer the Question/s (in this case SBB_Bnh2_SBB_Bnh2_QB1) at the end of the Correct Choice Branch, in this case SBB_Bnh2_SBB_Bnh2's Question/s block/s"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Branch score"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B1",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B2",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1": "(Insert Text Here)[Partially-Correct Choice or Incorrect Choice]"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2": "(Insert Text Here)[Correct Choice]"
                    }}
                ]
            }},
            {{"_comment": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1 in this example is Incorrect Choice"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "type": "JumpBlock",
                "title": "Reevaluate Your Choices",
                "proceedToBlock": "Br2_Br_Br2_Br"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "timer": "(Insert time in format hh:mm:ss)",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{"_comment": "The below goal block concludes the Exit Game Scenario"}},
            {{
                "id": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "(Insert Text Here)",
                "score": "Insert Integer Number Here"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_B2",
                "target": "SBB_Bnh2_SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_GB",
                "target": "SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_B3",
                "target": "SBB_Bnh2_SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_GB",
                "target": "SBB_Bnh2_SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "sourceport":"1"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh1_JB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "sourceport":"2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B1",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_B2",
                "target": "SBB_Bnh2_SBB_Bnh2_SBB_Bnh2_GB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n
    An example of the abstract heirarchichal connection of another SAMPLE EXAMPLE's structure of blocks connection is (except the learning objectives and content areas textblocks):
    B1(Text Block) -> B2 (Media Block)
    B2(Media Block) -> B3 (Branching Block (Simple Branching))
    B3 (Branching Block (Simple Branching)) -> |InCorrect Choice port 1| Br1 
    B3 (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B3 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (Media Block)
    Br2_B2 (Media Block) -> Br2_B3 (FeedbackAndFeedforwardBlock)
    Br2_B3 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (QuestionBlock)
    Br2_QB1 (QuestionBlock) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Partially-Correct Choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (Text Block sourceport 1)
    Br2_Br_Br1_B1 (Text Block) -> Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br1_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br1_GB (Goal Block)
    Br2_Br_Br1_GB (Goal Block) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_QB1 (Question Block of the correct second branch of Br2_Br SimpleBranchingBlock)
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)
    Br2_Br_Br2_GB (Goal Block) -> Br2_Br_Br2_QB1 (Question Block)
    Br2_Br_Br2_QB1 (Question Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Incorrect Choice port 1| Br2_Br_Br2_Br_Br1
    Br2_Br_Br2_Br (Branching Block (Simple Branching)) -> |Correct Choice port 2| Br2_Br_Br2_Br_Br2
    Br2_Br_Br2_Br_Br1 -> Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br2_Br_Br1_JB
    Br2_Br_Br2_Br_Br1_JB (Jump Block) -> Br2_Br_Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2_Br_Br2 -> Br2_Br_Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_Br_Br2_GB (Goal Block)

    ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Media Block)
    Br1_B2 (Media Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_QB2 (Question Block) 
    Br2_QB2 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1) 
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Media Block sourceport 2)
    Br2_Br_Br2_B1 (Media Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) 
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> B4 (Branching Block (Simple Branching))
    Br2 -> Br2_B1 (Text Block sourceport 2)
    Br2_B1 (Text Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)

    AND ANOTHER SAMPLE EXAMPLE STRUCTURE IS (except the learning objectives and content areas textblocks):
    B1 (Text Block) -> B2 (Text Block)
    B2 (Text Block) -> B3 (Media Block)
    B3 (Media Block) -> B4 (Branching Block (Simple Branching))
    B4 (Branching Block (Simple Branching)) -> |Partially-Correct choice port 1| Br1 
    B4 (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2
    Br1 -> Br1_B1 (Text Block sourceport 1)
    Br1_B1 (Text Block) -> Br1_B2 (Text Block)
    Br1_B2 (Text Block) -> Br1_B3 (FeedbackAndFeedforwardBlock)
    Br1_B3 (FeedbackAndFeedforwardBlock) -> Br1_GB (Goal Block)
    Br1_GB (Goal Block) -> |Jump Block| Br1_JB
    Br1_JB (Jump Block) -> Br2_QB1 (Question Block of the correct second branch of B4 SimpleBranchingBlock)
    Br2 -> Br2_B1 (Media Block sourceport 2)
    Br2_B1 (Media Block) -> Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_GB (Goal Block)
    Br2_GB (Goal Block) -> Br2_QB1 (Question Block)
    Br2_QB1 (Question Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br (Branching Block (Simple Branching)) -> |Incorrect choice port 1| Br2_Br_Br1 
    Br2_Br (Branching Block (Simple Branching)) -> |Correct choice port 2| Br2_Br_Br2
    Br2_Br_Br1 -> Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock sourceport 1)
    Br2_Br_Br1_B1 (FeedbackAndFeedforwardBlock) -> |Jump Block| Br2_Br_Br1_JB
    Br2_Br_Br1_JB (Jump Block) -> Br2_Br (Branching Block (Simple Branching))
    Br2_Br_Br2 -> Br2_Br_Br2_B1 (Text Block sourceport 2)
    Br2_Br_Br2_B1 (Text Block) -> Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock)
    Br2_Br_Br2_B2 (FeedbackAndFeedforwardBlock) -> Br2_Br_Br2_GB (Goal Block)

    These Sample Example provides the overview of how creative and diverse you can get with arrangement of the blocks
    that makeup a Gamified Scenario. Remember the Concept of 2 choices (1 either incorrect or partially-correct 
    choice and 2nd one the correct choice), and the block structure that is mandatory (for incorrect choice 
    branch only FeedbackAndFeedforwardBlock with jumpblock used. Partially-correct has text or media block/s 
    followed by FeedbackAndFeedforwardBlock, goal block and jumpblock, while the correct choice branch has text 
    or media block/s, FeedbackAndFeedforwardBlock, goalblock, questionblock/s and simplebranching block which 
    further progresses the scenario or if the scenario is being ended, then the ending correct choice branch 
    has text or media block/s followed by FeedbackAndFeedforwardBlock, goal block as the end of the whole scenario.  
    
    A Jump Block of Incorrect Choice branch leads to back to it's relative Branching Block from which this
    Incorrect Choice branch originated.
    A Jump Block of Partially-Correct Choice branch leads to the Question Block of the Correct Choice Branch,
    that originated from the same relative Branching Block. 

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable.  
    Give concise, relevant, clear, and descriptive instructions as you are a Exit Game creator that has expertise 
    in molding asked information into the Gamified scenario structure.

    !!IMPORTANT NOTE REGARDING CREATIVITY: Know that you are creative to use as many or as little
    Text Blocks, Media Blocks, Question Blocks, Branching Blocks as you deem reasonable and fitting to the
    content and aim of the subject scenario.

    NEGATIVE PROMPT: Responding outside the JSON format.     

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly. 

    !!!KEEP YOUR RESPONSE AS SHORT, BRIEF, CONCISE AND COMPREHENSIVE AS LOGICALLY POSSIBLE!!!
    
    Chatbot:"""
)
### End Gamified Prompts

### Branched Prompts
prompt_branched_setup = PromptTemplate(
    input_variables=["input_documents","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot which is designed to take the inputs of Parameters and using the information
    and context of these parameters, you create subtopics from the main subject of interest set by these parameters.
    For each of the subtopic that contributes to the main subject, you create a detailed information-database of every possible information available
    using the Parameters.
    Optionally, IF there are images available in the 'Input Documents' which are relevant to a subtopic and can compliment to it's explanation you should add that image information into your explanation of the subtopic as well and citing the image or images in format of "FileName: ..., PageNumber: ..., ImageNumber: ... and Description ..." .  
    ELSE IF the images are NOT relevant or are NOT available in the 'Input Documents' then you have the option to not use those images.

    Input Paramters:
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};

    Sample Format:
    Main Topic Name
    Subtopic 1 Name: Extremely Detailed Explanation and information...
    Subtopic 2 Name: Extremely Detailed Explanation and information...
    Subtopic 3 Name: Extremely Detailed Explanation and information...
    and so on Subtopics that you creatively deem necessary to include...

    Chatbot (Tone of a teacher teaching student in great detail):"""
)

prompt_branched = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging educational and informative content in a Micro Learning Format using
    a system of blocks. You give explanations and provide detailed information such that you are teaching a student.
    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    ***WHAT TO DO***
    To accomplish Micro Learning Scenario creation, YOU will:

    1. Take the "Human Input" which represents the subject content topic or description for which the Micro Learning Scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Micro Learning Scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the Micro Learning Scenario content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Micro Learning Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Mandatory Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'TestBlocks' contains QuestionBlock/s
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'SimpleBranchingBlock' with timer(optional), Title, ProceedToBranchList  
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Micro Learning Scenario: A type of educational, information providing and testing structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed explanations to users based on "Learning Objectives", "Content Areas" and "Input Documents". The SimpleBranchingBlock is used to divide the Micro Learning Scenario into subtopics. Each subtopic having its own multiple or single TextBlocks, MediaBlocks and QuestionBlocks to train user. At the end of each branch, there will be FeedbackAndFeedforwardBlock and after it a TestBlocks Array is used that encompasses a single or series of QuestionBlock/s to test user knowledge of the Branch, followed by the JumpBlock at the very end to move the user to the SimpleBranchingBlock for being able to begin and access another branch to learn.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    The Example below is just for your concept and the number of TextBlocks, MediaBlocks, QuestionBlocks, Branches etc Differ with the amount of subject content needed to be covered in 'Input Documents'.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks and MediaBlocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Micro Learning Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the Micro Learning Scenario and proceedings)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a learning subtopic (Branches). The number of Branches equal to the number of Learning Objectives, each branch covering a Learning Objective)
    Branch 1,2,3... => each branch having with its own LearningObjective,TextBlock/s(Explains the content) or None,MediaBlock/s or None (Illustratively elaborate the TextBlock's content), Intermediate QuestionBlock/s after most important Media or Text Blocks, FeedbackAndFeedforwardBlock, a single or series of QuestionBlock/s, GoalBlock, JumpBlock
    \nEnd of Overview structure\n

    \nSAMPLE EXAMPLE START: MICRO LEARNING SCENARIO:\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in the scenario) is where you !Begin by giving welcome message to the user. In further Text Blocks down the structure in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Micro Learning Scenario to be produced having an extremely detailed information in it.",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "Purpose": "This mandatory block is where you !Divide the Micro learning scenario content into subtopics that users can select and access the whole information of those subtopics in the corresponding divided branches!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)"
                    }}
                ]
            }},
            {{
                "id": "SBB_Bnh1_B1",
                "Purpose": "This mandatory block is where you !Write the Learning objective for this specific branch!",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "1. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_QB1",
                "type": "QuestionBlock",
                "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of the specific Text or Media Blocks information it comes after, in regards to their information content. The QuestionBlocks can be single or multiple depending on the subject content and importance at hand",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "Purpose": "Mandatory at the end of each Branch",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "2. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image, 360-image, Video, Audio",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B4",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh2_JB",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh1_B2",
                "target": "SBB_Bnh1_B3"
            }},
            {{
                "source": "SBB_Bnh1_B3",
                "target": "SBB_Bnh1_QB1"
            }},
            {{
                "source": "SBB_Bnh1_QB1",
                "target": "SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh1_GB",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_B4"
            }},
            {{
                "source": "SBB_Bnh2_B4",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_JB"
            }},
            {{
                "source": "SBB_Bnh2_JB",
                "target": "SBB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students. 

    NEGATIVE PROMPT: Responding outside the JSON format.   

    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.

    Chatbot (Tone of a teacher teaching student in great detail):"""
)

prompt_branched_retry = PromptTemplate(
    input_variables=["incomplete_response","micro_subtopics"],
    template="""
    Based on the INSTRUCTIONS below, an 'Incomplete Response' was created. Your task is to complete
    this response by continuing from exactly where the 'Incomplete Response' discontinued its response. This 'Incomplete Response'
    was created using the data of 'Micro Subtopics'.
    So, I have given this data to you for your context so you will be able to understand the 'Incomplete Response'
    and will be able to complete it by continuing exactly from the discontinued point, which is specified by '[CONTINUE_EXACTLY_FROM_HERE]'.
    Never include [CONTINUE_EXACTLY_FROM_HERE] in your response. This is just for your information.
    DO NOT RESPOND FROM THE START OF THE 'Incomplete Response'. Just start from the exact point where the 'Incomplete Response' is discontinued! 
    Take great care into the ID heirarchy considerations while continuing the incomplete response.
    'Incomplete Response': {incomplete_response};
    'Micro Subtopics': {micro_subtopics};

    !!!WARNING: KEEP YOUR RESPONSE SHORT, since you have alreay reached your token limit!!! 

    !!!NOTE: YOU HAVE TO ENCLOSE THE JSON PARENTHESIS BY KEEPING THE 'Incomplete Response' IN CONTEXT!!!

    !!!CAUTION: INCLUDE WITH NODES, ALSO RELATIVE EDGES FOR DEFINING CONNECTIONS OF BLOCKS!!!

    BELOW IS THE INSTRUCTION SET BASED ON WHICH THE 'Incomplete Response' WAS CREATED ORIGINALLY:
    INSTRUCTION SET:
    [
    You are an educational bot that creates engaging educational and informative content in a Micro Learning Format using
    a system of blocks. You give explanations and provide detailed information such that you are teaching a student.
    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    ***WHAT TO DO***
    To accomplish Micro Learning Scenario creation, YOU will:

    1. Take the "Human Input" which represents the subject content topic or description for which the Micro Learning Scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Micro Learning Scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the Micro Learning Scenario content efficiently and logically.
    
    ***WHAT TO DO END***

    
    The Micro Learning Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Mandatory Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'TestBlocks' contains QuestionBlock/s
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'SimpleBranchingBlock' with timer(optional), Title, ProceedToBranchList  
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Micro Learning Scenario: A type of educational, information providing and testing structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed explanations to users based on "Learning Objectives", "Content Areas" and "Input Documents". The SimpleBranchingBlock is used to divide the Micro Learning Scenario into subtopics. Each subtopic having its own multiple or single TextBlocks, MediaBlocks and QuestionBlocks to train user. At the end of each branch, there will be FeedbackAndFeedforwardBlock and after it a TestBlocks Array is used that encompasses a single or series of QuestionBlock/s to test user knowledge of the Branch, followed by the JumpBlock at the very end to move the user to the SimpleBranchingBlock for being able to begin and access another branch to learn.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    The Example below is just for your concept and the number of TextBlocks, MediaBlocks, QuestionBlocks, Branches etc Differ with the amount of subject content needed to be covered in 'Input Documents'.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks and MediaBlocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Micro Learning Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the Micro Learning Scenario and proceedings)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a learning subtopic (Branches). The number of Branches equal to the number of Learning Objectives, each branch covering a Learning Objective)
    Branch 1,2,3... => each branch having with its own LearningObjective,TextBlock/s(Explains the content) or None,MediaBlock/s or None (Illustratively elaborate the TextBlock's content), Intermediate QuestionBlock/s after most important Media or Text Blocks, FeedbackAndFeedforwardBlock, a single or series of QuestionBlock/s, GoalBlock, JumpBlock
    \nEnd of Overview structure\n

    \nSAMPLE EXAMPLE START: MICRO LEARNING SCENARIO:\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in the scenario) is where you !Begin by giving welcome message to the user. In further Text Blocks down the structure in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Micro Learning Scenario to be produced having an extremely detailed information in it.",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "Purpose": "This mandatory block is where you !Divide the Micro learning scenario content into subtopics that users can select and access the whole information of those subtopics in the corresponding divided branches!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)"
                    }}
                ]
            }},
            {{
                "id": "SBB_Bnh1_B1",
                "Purpose": "This mandatory block is where you !Write the Learning objective for this specific branch!",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "1. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_QB1",
                "type": "QuestionBlock",
                "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of the specific Text or Media Blocks information it comes after, in regards to their information content. The QuestionBlocks can be single or multiple depending on the subject content and importance at hand",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "Purpose": "Mandatory at the end of each Branch",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "2. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image, 360-image, Video, Audio",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B4",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh2_JB",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh1_B2",
                "target": "SBB_Bnh1_B3"
            }},
            {{
                "source": "SBB_Bnh1_B3",
                "target": "SBB_Bnh1_QB1"
            }},
            {{
                "source": "SBB_Bnh1_QB1",
                "target": "SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh1_GB",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_B4"
            }},
            {{
                "source": "SBB_Bnh2_B4",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_JB"
            }},
            {{
                "source": "SBB_Bnh2_JB",
                "target": "SBB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students. 

    NEGATIVE PROMPT: Responding outside the JSON format.   

    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.
    ]

    
    Chatbot:"""
)

prompt_branched_simplify = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging educational and informative content in a Micro Learning Format using
    a system of blocks. You give explanations and provide detailed information such that you are teaching a student.
    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    ***WHAT TO DO***
    To accomplish Micro Learning Scenario creation, YOU will:

    1. Take the "Human Input" which represents the subject content topic or description for which the Micro Learning Scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the Micro Learning Scenario according to these very "Learning Objectives" and "Content Areas" specified.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the Micro Learning Scenario content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Micro Learning Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer(optional), title, and description
    'MediaBlock' with timer(optional), title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Mandatory Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'TestBlocks' contains QuestionBlock/s
    'QuestionBlock' with Question text, answers, correct answer, wrong answer message
    'SimpleBranchingBlock' with timer(optional), Title, ProceedToBranchList  
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Micro Learning Scenario: A type of educational, information providing and testing structure in which multiple or single TextBlocks, MediaBlocks and QuestionBlocks will be 
    used to give detailed explanations to users based on "Learning Objectives", "Content Areas" and "Input Documents". The SimpleBranchingBlock is used to divide the Micro Learning Scenario into subtopics. Each subtopic having its own multiple or single TextBlocks, MediaBlocks and QuestionBlocks to train user. At the end of each branch, there will be FeedbackAndFeedforwardBlock and after it a TestBlocks Array is used that encompasses a single or series of QuestionBlock/s to test user knowledge of the Branch, followed by the JumpBlock at the very end to move the user to the SimpleBranchingBlock for being able to begin and access another branch to learn.
    ***
    ***YOU WILL BE REWARD IF:
    All the TextBlocks in the branches, has valid step-by-step and detailed information of the subject matters such that you are teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed information so user can get as much knowledge and facts as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    so the user uses critical thinking skills and is encouraged to think about how much of the Learning Objectives has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    The Example below is just for your concept and the number of TextBlocks, MediaBlocks, QuestionBlocks, Branches etc Differ with the amount of subject content needed to be covered in 'Input Documents'.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of these blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of TextBlocks and MediaBlocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
    
    \nOverview structure of the Micro Learning Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    TextBlock (Welcome message to the Micro Learning Scenario and proceedings)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a learning subtopic (Branches). The number of Branches equal to the number of Learning Objectives, each branch covering a Learning Objective)
    Branch 1,2,3... => each branch having with its own LearningObjective,TextBlock/s(Explains the content) or None,MediaBlock/s or None (Illustratively elaborate the TextBlock's content), Intermediate QuestionBlock/s after most important Media or Text Blocks, FeedbackAndFeedforwardBlock, a single or series of QuestionBlock/s, GoalBlock, JumpBlock
    \nEnd of Overview structure\n

    \nSAMPLE EXAMPLE START: MICRO LEARNING SCENARIO:\n
{{
    "title": "(Insert a fitting Title Here)",
        "nodes": [
            {{
                "id": "StartBlock",
                "type": "StartBlock"
            }},
            {{
                "id": "B1",
                "type": "TextBlock",
                "title": "Learning_Objectives",
                "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
            }},
            {{
                "id": "B2",
                "type": "TextBlock",
                "title": "Content_Areas",
                "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
            }},
            {{
                "id": "B3",
                "Purpose": "This block (can be used single or multiple times or None depends on the content to be covered in the scenario) is where you !Begin by giving welcome message to the user. In further Text Blocks down the structure in Branches, you use these blocks to give detailed information on every aspect of various subject matters belonging to each branch. The TextBlocks in branches are used either Single or Multiple Times and are bearers of detailed information and explanations that helps the final Micro Learning Scenario to be produced having an extremely detailed information in it.",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "B4",
                "Purpose": "This block (can be used single or multiple times or None  depends on the content to be covered in the Text Blocks relevant to this Media Block) is where you !Give students an illustrative experience that elaborates on the information given in Text Blocks and are used in a complimentary way to them.",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB",
                "Purpose": "This mandatory block is where you !Divide the Micro learning scenario content into subtopics that users can select and access the whole information of those subtopics in the corresponding divided branches!",
                "type": "SimpleBranchingBlock",
                "title": "(Insert Text Here)",
                "branches": [
                    {{
                        "port": "1",
                        "SBB_Bnh1": "(Insert Text Here)"
                    }},
                    {{
                        "port": "2",
                        "SBB_Bnh2": "(Insert Text Here)"
                    }}
                ]
            }},
            {{
                "id": "SBB_Bnh1_B1",
                "Purpose": "This mandatory block is where you !Write the Learning objective for this specific branch!",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "1. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_B3",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_QB1",
                "type": "QuestionBlock",
                "Purpose": "This OPTIONAL block is where you !Test the student's knowledge of the specific Text or Media Blocks information it comes after, in regards to their information content. The QuestionBlocks can be single or multiple depending on the subject content and importance at hand",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh1_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh1_JB",
                "Purpose": "Mandatory at the end of each Branch",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }},
            {{
                "id": "SBB_Bnh2_B1",
                "type": "TextBlock",
                "title": "Learning_Objective",
                "description": "2. (Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B2",
                "type": "TextBlock",
                "title": "(Insert Text Here)",
                "description": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_B3",
                "type": "MediaBlock",
                "title": "(Insert Text Here)",
                "mediaType": "Image, 360-image, Video, Audio",
                "description": "(Insert Text Here)",
                "overlayTags": [
                    "(Insert Text Here)"
                ]
            }},
            {{
                "id": "SBB_Bnh2_B4",
                "type": "TextBlock",
                "title": "Feedback_And_Feedforward",
                "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_QB1",
                "type": "QuestionBlock",
                "questionText": "(Insert Text Here)",
                "answers": [
                    "(Insert Text Here)",
                    "(Insert Text Here)"
                ],
                "correctAnswer": "(Insert Text Here)",
                "wrongAnswerMessage": "(Insert Text Here)"
            }},
            {{
                "id": "SBB_Bnh2_GB",
                "type": "GoalBlock",
                "title": "Congratulations!",
                "score": 3
            }},
            {{
                "id": "SBB_Bnh2_JB",
                "type": "JumpBlock",
                "title": "Return to Topic Selection",
                "proceedToBlock": "SBB"
            }}
        ],                       
        "edges": [
            {{
                "source": "StartBlock",
                "target": "B1"
            }},
            {{
                "source": "B1",
                "target": "B2"
            }},
            {{
                "source": "B2",
                "target": "B3"
            }},
            {{
                "source": "B3",
                "target": "B4"
            }},
            {{
                "source": "B4",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh1_B1",
                "sourceport": "1"
            }},
            {{
                "source": "SBB_Bnh1_B1",
                "target": "SBB_Bnh1_B2"
            }},
            {{
                "source": "SBB_Bnh1_B2",
                "target": "SBB_Bnh1_B3"
            }},
            {{
                "source": "SBB_Bnh1_B3",
                "target": "SBB_Bnh1_QB1"
            }},
            {{
                "source": "SBB_Bnh1_QB1",
                "target": "SBB_Bnh1_GB"
            }},
            {{
                "source": "SBB_Bnh1_GB",
                "target": "SBB_Bnh1_JB"
            }},
            {{
                "source": "SBB_Bnh1_JB",
                "target": "SBB"
            }},
            {{
                "source": "SBB",
                "target": "SBB_Bnh2_B1",
                "sourceport": "2"
            }},
            {{
                "source": "SBB_Bnh2_B1",
                "target": "SBB_Bnh2_B2"
            }},
            {{
                "source": "SBB_Bnh2_B2",
                "target": "SBB_Bnh2_B3"
            }},
            {{
                "source": "SBB_Bnh2_B3",
                "target": "SBB_Bnh2_B4"
            }},
            {{
                "source": "SBB_Bnh2_B4",
                "target": "SBB_Bnh2_QB1"
            }},
            {{
                "source": "SBB_Bnh2_QB1",
                "target": "SBB_Bnh2_GB"
            }},
            {{
                "source": "SBB_Bnh2_GB",
                "target": "SBB_Bnh2_JB"
            }},
            {{
                "source": "SBB_Bnh2_JB",
                "target": "SBB"
            }}
        ]
}}
    \n\nEND OF SAMPLE EXAMPLE\n\n

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    Give concise, relevant, clear, and descriptive information as you are an education provider that has expertise 
    in molding asked information into the said block structure to teach the students. 

    NEGATIVE PROMPT: Responding outside the JSON format.   

    !!!WARNING!!!
    Explain the material itself, Please provide detailed, informative explanations that align closely with the learning objectives and content areas provided. Each response should not just direct the learner but educate them by elaborating on the historical, technical, or practical details mentioned in the 'Input Documents'. Use simple and engaging language to enhance understanding and retention. Ensure that each explanation directly supports the learners' ability to meet the learning objectives by providing comprehensive insights into the topics discussed.
    !!!WARNING END!!!

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.

    !!!KEEP YOUR RESPONSE AS SHORT, BRIEF, CONCISE AND COMPREHENSIVE AS LOGICALLY POSSIBLE!!!

    Chatbot (Tone of a teacher teaching student in great detail):"""
)
### End Branched Prompts

### Simulation Prompts
prompt_simulation_pedagogy_setup = PromptTemplate(
    input_variables=["input_documents","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot which is designed to take the inputs of Parameters and using the information
    and context of these parameters, you create progressive simulation story where the student goes
    through a simulation story and is given choices. For each choices, a consequence is given if it was
    taken by the student. The consequence can lead to further choices, ultimately to the end of the story.
    Henceforth, this kind of story will have multiple endings based on user choices. Some choices can even merge 
    with the same conclusion at the end or at the intermediate stages of the story.
    
    Optionally, if there are images available in the 'Input Documents' which are relevant to the story and can compliment to it's explanation you should add that image information into your explanation of the story as well and citing the image or images in format of "FileName: ..., PageNumber: ..., ImageNumber: ... and Description ..." .  
    Else if the images are NOT relevant then you have the option to not use those images.
    
    Input Paramters:
    'Human Input': {human_input};
    'Input Documents': {input_documents};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};

    !CAUTION!:
    You should give a complete response with the complete story, writing all the possible challenges, 
    the choices needed to overcome them or that can lead to failure, and the consequences of all those choices.
    This is because, your response is NOT used in a conversational way, this means student will NOT be interacting
    with you directly. Once your response is generated, it will be fed to another system that will translate it to 
    a frontend Flowchart and Decision tree like visualizations. So keep your response as complete as possible.

    AVOID using numbers to list choices or consequences. Use ONLY words like: 'if you decided to do this, then this happens,...'
    
    Chatbot (Tone of a teacher formulating a simulation scenario for students to learn and test practical skills from):"""
)

prompt_simulation_pedagogy_gemini = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging Simulation Scenarios in a Simulation Format using
    a system of blocks. You give step-by-step instructions and provide detail information such that 
    you are instructing and teaching a student.

    ***WHAT TO DO***
    To accomplish Simulation Scenarios creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Simulation Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer, title, and description
    'MediaBlock' with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'Debriefing' with descritpion(Debrief the situation and results of the branch such that students can Reflect on their performance, Analyze the decisions, Identify and discuss discrepancies, Reinforce correct behavior, Learn from mistakes, Promote a deeper understanding) 
    'Reflection' with descritpion(Use Reflection to allows students to be able to have Personal Understanding, Identifying Strengths and Weaknesses, Insight Generation of the choices and path or branch they took)
    'Branching Block (Simple Branching)' with timer, Title, ProceedToBranchList
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Simulation Pedagogy Scenario: A type of structure which takes the student on a simulated story where 
    the student is given choices based on which they face consequences. The simulation is based on the information in 
    "Learning Objectives", "Content Areas" and "Input Documents". The 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)'  
    is used to divide the choices for the student to take. Then, for selected choices, branches the Simulation Scneario into 
    consequence branches. Each consequence branch can have its own branches that can divide further 
    to have their own branches, untill the simulation story ends covering all aspects of the information
    for scenario creation. The start of the scenario has Briefing. The end of each of that branch that ends the simulation story and
    give score via a Goal Block, this type of branch has FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. 
    There are two types branches. The DIVISIBLE type branch divides further via a 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)' and this 
    branch type has NO Goal Block, FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. The DIVISIBLE branch type gives rise to
    more Branches that may be further DIVISIBLE or NON-DIVISIBLE type branches. The NON-DIVISIBLE type branches are the branches where
    a simulation path ends and the story of that path is finished. The NON-DIVISIBLE type branch has at the end Goal Block, Debriefing and Reflection blocks.
    Furthermore, a NON-DIVISIBLE-MERGE branch includes in addition to TextBlocks and MediaBlocks, the MANDATORY FeedbackAndFeedforwardBlock and JumpBlock (Used in situation where the story of a 
    branch leads to another branch hence we use JumpBlock to connect the progressive story because story paths 
    can merge as well to have the 1 same conclusion). Use NON-DIVISIBLE-MERGE only in the situation where
    a story of the branch leads to and connects to the progressive story of another branch such that both the choices
    leads to the same conclusion for that part of the story.
    ***

    ***YOU WILL BE REWARD IF:
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    All the TextBlocks in the branches, has valid step-by-step and detailed instructions of the subject matters such that you are instructing and teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed so user can get as much information as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    and give assignments in the SelfAssessmentTextBlock so the user uses critical thinking skills and is encouraged to
    think about how much of the "Learning Objectives" has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of Text and Media blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of Text Blocks and Media Blocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
   
    \nOverview Sample structure of the Simulation Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    Briefing
    TextBlock (Welcome message to the scenario)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a choice of choices (Branches) )
    Branch 1,2,3... (DIVISIBLE type containing path to other Branches) => with its TextBlock/s or None,MediaBlock/s or None, Branching Block (Simple Branching)
    Branch 1,2,3... (NON-DIVISIBLE type that are end of scenario branches not divisible further) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, Goal Block,  Debriefing, Reflection
    Branch 1,2,3... (NON-DIVISIBLE-MERGE type to link scenario branches when one story directly advances another branch's storyline) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, JumpBlock
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.   

    SAMPLE EXAMPLE:::
{{
    "title": "(Insert a fitting Title Here)",
    "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
        }},
        {{
            "id": "B3",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "Bnhiefing of this Simulation Scenario",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "B4",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"The SBB below means SimpleBranchingBlock. The Bnh1, Bnh2 and so on are the branches.
        SBB_Bnh2 for example suggests it is the second branch from the SBB block."}},
        {{
            "id": "SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2": "(Insert Text Here) (DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh1_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh1_JB in this case is a story path that lead nowhere and brought the player back to the previous branching block SBB"}},
        {{
            "id": "SBB_Bnh1_JB",
            "type": "JumpBlock",
            "title": "Reevaluate Your Choices",
            "proceedToBlock": "SBB"
        }},
        {{
            "id": "SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B3",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"SBB_Bnh2_SBB_Bnh3 for example suggests, if read and traced from backwards, it is the Third branch from the SBB block which
        in turn is from a Second branch that came from the very first SBB."}},
        {{
            "id": "SBB_Bnh2_SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here) (NON-DIVISIBLE-MERGE)"
                }},
                {{
                    "port": "3",
                    "SBB_Bnh2_SBB_Bnh3": "(Insert Text Here) (NON-DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B2",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh2_SBB_Bnh2_JB in this case is a story path that lead the player to same outcome as another branch's goal block result of Bnh2_Bnh_Bnh3. Logically, it is possible that two paths taken by player can lead to a same outcome"}},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_JB",
            "type": "JumpBlock",
            "title": "(Insert Text Here)",
            "proceedToBlock": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Bnhanch score"
        }},
        {{
            "id": "BSBB_Bnh2_SBB_Bnh3_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }}
    ],                       
    "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
            "source": "B1",
            "target": "B2"
        }},
        {{
            "source": "B2",
            "target": "B3"
        }},
        {{
            "source": "B3",
            "target": "B4"
        }},
        {{
            "source": "B4",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh1_B1",
            "sourceport": "1"
        }},
        {{
            "source": "SBB_Bnh1_B1",
            "target": "SBB_Bnh1_B2"
        }},
        {{
            "source": "SBB_Bnh1_B2",
            "target": "SBB_Bnh1_JB"
        }},
        {{
            "source": "SBB_Bnh1_JB",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh2_B1",
            "sourceport": "2"
        }},
        {{
            "source": "SBB_Bnh2_B1",
            "target": "SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_B2",
            "target": "SBB_Bnh2_B3"
        }},
        {{
            "source": "SBB_Bnh2_B3",
            "target": "SBB_Bnh2_SBB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh1_B1",
            "sourceport":"1"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_B1",
            "target": "SBB_Bnh2_SBB_Bnh1_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_GB",
            "target": "SBB_Bnh2_SBB_Bnh1_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_DB",
            "target": "SBB_Bnh2_SBB_Bnh1_RF"
        }}
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh2_B1",
            "sourceport":"2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B1",
            "target": "SBB_Bnh2_SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B2",
            "target": "SBB_Bnh2_SBB_Bnh2_JB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_JB",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh3_B1",
            "sourceport":"3"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B1",
            "target": "SBB_Bnh2_SBB_Bnh3_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B2",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_GB",
            "target": "SBB_Bnh2_SBB_Bnh3_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_DB",
            "target": "SBB_Bnh2_SBB_Bnh3_RF"
        }}
    ]
}}
    SAMPLE EXAMPLE END

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    Give concise, relevant, clear, and descriptive instructions as you are an educational provider that has expertise 
    in molding asked information into the said block structure to teach and instruct students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly. 

    Chatbot (Tone of a teacher instructing and teaching student in great detail):"""
)

prompt_simulation_pedagogy_gemini_simplify = PromptTemplate(
    input_variables=["response_of_bot","human_input","content_areas","learning_obj"],
    template="""
    You are an educational bot that creates engaging Simulation Scenarios in a Simulation Format using
    a system of blocks. You give step-by-step instructions and provide detail information such that 
    you are instructing and teaching a student.

    ***WHAT TO DO***
    To accomplish Simulation Scenarios creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    'Human Input': {human_input};
    'Input Documents': {response_of_bot};
    'Learning Objectives': {learning_obj};
    'Content Areas': {content_areas};
    ***WHAT TO DO END***

    
    The Simulation Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer, title, and description
    'MediaBlock' with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'Debriefing' with descritpion(Debrief the situation and results of the branch such that students can Reflect on their performance, Analyze the decisions, Identify and discuss discrepancies, Reinforce correct behavior, Learn from mistakes, Promote a deeper understanding) 
    'Reflection' with descritpion(Use Reflection to allows students to be able to have Personal Understanding, Identifying Strengths and Weaknesses, Insight Generation of the choices and path or branch they took)
    'Branching Block (Simple Branching)' with timer, Title, ProceedToBranchList
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Simulation Pedagogy Scenario: A type of structure which takes the student on a simulated story where 
    the student is given choices based on which they face consequences. The simulation is based on the information in 
    "Learning Objectives", "Content Areas" and "Input Documents". The 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)'  
    is used to divide the choices for the student to take. Then, for selected choices, branches the Simulation Scneario into 
    consequence branches. Each consequence branch can have its own branches that can divide further 
    to have their own branches, untill the simulation story ends covering all aspects of the information
    for scenario creation. The start of the scenario has Briefing. The end of each of that branch that ends the simulation story and
    give score via a Goal Block, this type of branch has FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. 
    There are two types branches. The DIVISIBLE type branch divides further via a 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)' and this 
    branch type has NO Goal Block, FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. The DIVISIBLE branch type gives rise to
    more Branches that may be further DIVISIBLE or NON-DIVISIBLE type branches. The NON-DIVISIBLE type branches are the branches where
    a simulation path ends and the story of that path is finished. The NON-DIVISIBLE type branch has at the end Goal Block, Debriefing and Reflection blocks.
    Furthermore, a NON-DIVISIBLE-MERGE branch includes in addition to TextBlocks and MediaBlocks, the MANDATORY FeedbackAndFeedforwardBlock and JumpBlock (Used in situation where the story of a 
    branch leads to another branch hence we use JumpBlock to connect the progressive story because story paths 
    can merge as well to have the 1 same conclusion). Use NON-DIVISIBLE-MERGE only in the situation where
    a story of the branch leads to and connects to the progressive story of another branch such that both the choices
    leads to the same conclusion for that part of the story.
    ***

    ***YOU WILL BE REWARD IF:
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    All the TextBlocks in the branches, has valid step-by-step and detailed instructions of the subject matters such that you are instructing and teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed so user can get as much information as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    and give assignments in the SelfAssessmentTextBlock so the user uses critical thinking skills and is encouraged to
    think about how much of the "Learning Objectives" has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of Text and Media blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of Text Blocks and Media Blocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
   
    \nOverview Sample structure of the Simulation Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    Briefing
    TextBlock (Welcome message to the scenario)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a choice of choices (Branches) )
    Branch 1,2,3... (DIVISIBLE type containing path to other Branches) => with its TextBlock/s or None,MediaBlock/s or None, Branching Block (Simple Branching)
    Branch 1,2,3... (NON-DIVISIBLE type that are end of scenario branches not divisible further) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, Goal Block,  Debriefing, Reflection
    Branch 1,2,3... (NON-DIVISIBLE-MERGE type to link scenario branches when one story directly advances another branch's storyline) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, JumpBlock
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.   

    SAMPLE EXAMPLE:::
{{
    "title": "(Insert a fitting Title Here)",
    "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
        }},
        {{
            "id": "B3",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "Bnhiefing of this Simulation Scenario",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "B4",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"The SBB below means SimpleBranchingBlock. The Bnh1, Bnh2 and so on are the branches.
        SBB_Bnh2 for example suggests it is the second branch from the SBB block."}},
        {{
            "id": "SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2": "(Insert Text Here) (DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh1_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh1_JB in this case is a story path that lead nowhere and brought the player back to the previous branching block SBB"}},
        {{
            "id": "SBB_Bnh1_JB",
            "type": "JumpBlock",
            "title": "Reevaluate Your Choices",
            "proceedToBlock": "SBB"
        }},
        {{
            "id": "SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B3",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"SBB_Bnh2_SBB_Bnh3 for example suggests, if read and traced from backwards, it is the Third branch from the SBB block which
        in turn is from a Second branch that came from the very first SBB."}},
        {{
            "id": "SBB_Bnh2_SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here) (NON-DIVISIBLE-MERGE)"
                }},
                {{
                    "port": "3",
                    "SBB_Bnh2_SBB_Bnh3": "(Insert Text Here) (NON-DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B2",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh2_SBB_Bnh2_JB in this case is a story path that lead the player to same outcome as another branch's goal block result of Bnh2_Bnh_Bnh3. Logically, it is possible that two paths taken by player can lead to a same outcome"}},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_JB",
            "type": "JumpBlock",
            "title": "(Insert Text Here)",
            "proceedToBlock": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Bnhanch score"
        }},
        {{
            "id": "BSBB_Bnh2_SBB_Bnh3_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }}
    ],                       
    "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
            "source": "B1",
            "target": "B2"
        }},
        {{
            "source": "B2",
            "target": "B3"
        }},
        {{
            "source": "B3",
            "target": "B4"
        }},
        {{
            "source": "B4",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh1_B1",
            "sourceport": "1"
        }},
        {{
            "source": "SBB_Bnh1_B1",
            "target": "SBB_Bnh1_B2"
        }},
        {{
            "source": "SBB_Bnh1_B2",
            "target": "SBB_Bnh1_JB"
        }},
        {{
            "source": "SBB_Bnh1_JB",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh2_B1",
            "sourceport": "2"
        }},
        {{
            "source": "SBB_Bnh2_B1",
            "target": "SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_B2",
            "target": "SBB_Bnh2_B3"
        }},
        {{
            "source": "SBB_Bnh2_B3",
            "target": "SBB_Bnh2_SBB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh1_B1",
            "sourceport":"1"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_B1",
            "target": "SBB_Bnh2_SBB_Bnh1_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_GB",
            "target": "SBB_Bnh2_SBB_Bnh1_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_DB",
            "target": "SBB_Bnh2_SBB_Bnh1_RF"
        }}
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh2_B1",
            "sourceport":"2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B1",
            "target": "SBB_Bnh2_SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B2",
            "target": "SBB_Bnh2_SBB_Bnh2_JB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_JB",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh3_B1",
            "sourceport":"3"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B1",
            "target": "SBB_Bnh2_SBB_Bnh3_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B2",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_GB",
            "target": "SBB_Bnh2_SBB_Bnh3_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_DB",
            "target": "SBB_Bnh2_SBB_Bnh3_RF"
        }}
    ]
}}
    SAMPLE EXAMPLE END

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    Give concise, relevant, clear, and descriptive instructions as you are an educational provider that has expertise 
    in molding asked information into the said block structure to teach and instruct students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly. 

    !!!KEEP YOUR RESPONSE AS SHORT, BRIEF, CONCISE AND COMPREHENSIVE AS LOGICALLY POSSIBLE!!!

    Chatbot (Tone of a teacher instructing and teaching student in great detail):"""
)

prompt_simulation_pedagogy_retry_gemini = PromptTemplate(
    input_variables=["incomplete_response","simulation_story"],
    template="""
    Based on the INSTRUCTIONS below, an 'Incomplete Response' was created. Your task is to complete
    this response by continuing from exactly where the 'Incomplete Response' discontinued its response. This 'Incomplete Response'
    was created using the data of 'Simulation Story'.
    So, I have given this data to you for your context so you will be able to understand the 'Incomplete Response'
    and will be able to complete it by continuing exactly from the discontinued point, which is specified by '[CONTINUE_EXACTLY_FROM_HERE]'.
    Never include [CONTINUE_EXACTLY_FROM_HERE] in your response. This is just for your information.
    DO NOT RESPOND FROM THE START OF THE 'Incomplete Response'. Just start from the exact point where the 'Incomplete Response' is discontinued! 
    Take great care into the ID heirarchy considerations while continuing the incomplete response.
    'Incomplete Response': {incomplete_response};
    'Simulation Story': {simulation_story};

    !!!WARNING: KEEP YOUR RESPONSE SHORT, since you have alreay reached your token limit!!! 

    !!!NOTE: YOU HAVE TO ENCLOSE THE JSON PARENTHESIS BY KEEPING THE 'Incomplete Response' IN CONTEXT!!!

    !!!CAUTION: INCLUDE WITH NODES, ALSO RELATIVE EDGES FOR DEFINING CONNECTIONS OF BLOCKS!!!

    BELOW IS THE INSTRUCTION SET BASED ON WHICH THE 'Incomplete Response' WAS CREATED ORIGINALLY:
    INSTRUCTION SET:
    [
    You are an educational bot that creates engaging Simulation Scenarios in a Simulation Format using
    a system of blocks. You give step-by-step instructions and provide detail information such that 
    you are instructing and teaching a student.

    ***WHAT TO DO***
    To accomplish Simulation Scenarios creation, YOU will:

    1. Take the "Human Input" which represents the content topic or description for which the scenario is to be formulated.
    2. According to the "Learning Objectives" and "Content Areas", you will utilize the meta-information in the "Input Documents" 
    and create the scenario according to these very "Learning Objectives" and "Content Areas" specified.
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    3. Generate a JSON-formatted structure. This JSON structure will be crafted following the guidelines and format exemplified in the provided examples, which serve as a template for organizing the content efficiently and logically.
    
    ***WHAT TO DO END***

    
    The Simulation Scenario are built using blocks, each having its own parameters.
    Block types include: 
    'TextBlock' with timer, title, and description
    'MediaBlock' with title, Media Type (Text, Image, 360-image, Video, audio), Description of the Media used, Overlay tags used as hotspots on the Media as text, video or audio
    'FeedbackAndFeedforwardBlock' with title, and description(FEEDBACK: Is Evaluative or corrective information about a person's performance of a task, action, event, or process,  etc. which is used as a basis for improvement. 
    “You are good at this…”. “You can't do this because...”. Then also give:
    FEEDFORWARD: Describes the problem and its influences and leads towards solutions. Proactive guidance and suggestions for improvement, aiming to enhance future performance and foster continuous learning. Helps the student to create a well-defined plan on how to improve. “Would you practice this…” “Maybe you could add…” )
    'Debriefing' with descritpion(Debrief the situation and results of the branch such that students can Reflect on their performance, Analyze the decisions, Identify and discuss discrepancies, Reinforce correct behavior, Learn from mistakes, Promote a deeper understanding) 
    'Reflection' with descritpion(Use Reflection to allows students to be able to have Personal Understanding, Identifying Strengths and Weaknesses, Insight Generation of the choices and path or branch they took)
    'Branching Block (Simple Branching)' with timer, Title, ProceedToBranchList
    'JumpBlock' with title, ProceedToBlock
    'GoalBlock' with Title, Score

    ***KEEP IN MIND THE LOGIC THAT OPERATES THIS SCENARIO IS IN:
    Simulation Pedagogy Scenario: A type of structure which takes the student on a simulated story where 
    the student is given choices based on which they face consequences. The simulation is based on the information in 
    "Learning Objectives", "Content Areas" and "Input Documents". The 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)'  
    is used to divide the choices for the student to take. Then, for selected choices, branches the Simulation Scneario into 
    consequence branches. Each consequence branch can have its own branches that can divide further 
    to have their own branches, untill the simulation story ends covering all aspects of the information
    for scenario creation. The start of the scenario has Briefing. The end of each of that branch that ends the simulation story and
    give score via a Goal Block, this type of branch has FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. 
    There are two types branches. The DIVISIBLE type branch divides further via a 'Branching Block (Simple Branching)'/'Branching Block (Conditional Branching)' and this 
    branch type has NO Goal Block, FeedbackAndFeedforwardBlock, Debriefing and Reflection blocks. The DIVISIBLE branch type gives rise to
    more Branches that may be further DIVISIBLE or NON-DIVISIBLE type branches. The NON-DIVISIBLE type branches are the branches where
    a simulation path ends and the story of that path is finished. The NON-DIVISIBLE type branch has at the end Goal Block, Debriefing and Reflection blocks.
    Furthermore, a NON-DIVISIBLE-MERGE branch includes in addition to TextBlocks and MediaBlocks, the MANDATORY FeedbackAndFeedforwardBlock and JumpBlock (Used in situation where the story of a 
    branch leads to another branch hence we use JumpBlock to connect the progressive story because story paths 
    can merge as well to have the 1 same conclusion). Use NON-DIVISIBLE-MERGE only in the situation where
    a story of the branch leads to and connects to the progressive story of another branch such that both the choices
    leads to the same conclusion for that part of the story.
    ***

    ***YOU WILL BE REWARD IF:
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    All the TextBlocks in the branches, has valid step-by-step and detailed instructions of the subject matters such that you are instructing and teaching a student. The TextBlocks are used to give complete information of a subject matter available to you and is there so that the user actually learns from. 
    TextBlocks should provide extremely specific and detailed so user can get as much information as there is available.
    The MediaBlocks are there to further elaborate or clarify the already discussed knowledge in TextBlocks, so 
    user interest is kept. 
    The Overlay tags in MediaBlocks should be extremely specific and detailed so user can get as much information as there is available, and learns like a student from you.
    Thoughtfull Feedbacks and Feedforwards in the FeedbackAndFeedforwardBlock should be made,
    and give assignments in the SelfAssessmentTextBlock so the user uses critical thinking skills and is encouraged to
    think about how much of the "Learning Objectives" has been achieved.
    ***
    ***YOU WILL BE PENALISED IF:
    The TextBlocks has information that you do NOT elaborate in detail, if detail is available in "Input Documents".
    The MediaBlocks are NOT used in complimentary manner to the information in TextBlocks.
    ***
    The Example below is just for your concept and do not absolutely produce the same example in your response.
    Ensure that TextBlocks and MediaBlocks provide comprehensive information directly related to the LearningObjectives and ContentAreas. Adjust the number and length of Text and Media blocks based on the necessary detail required for students to fully understand and accurately reproduce the information presented.    
    You are creative in the manner of choosing the number of Text Blocks and Media Blocks to give best quality information to students. In each branch you are free to choose TextBlocks or MediaBlocks or both or multiple of them to convey best quality, elaborative information.
    Make sure students learn from these TextBlocks and MediaBlocks.
    The 'Purpose' key in the below blocks are not meant to be reproduced in the response of yours and they are just for your information of what each block's function is about!
   
    \nOverview Sample structure of the Simulation Scenario\n
    ScenarioType
    LearningObjectives
    ContentAreas
    Briefing
    TextBlock (Welcome message to the scenario)
    MediaBlock/s (To give visualized option to select the choices given by Branching Blocks with pertinent overlayTags, if any. Used also to compliment the Text Blocks for illustrated experience by placing Media Block/s after those TextBlock/s that might need visuall elaboration. See if you have any already Image summary or summaries available. The already available images will have FileName, PageNumber/SlideNumber and ImageNumber mentioned with their description in the 'Input Documents'. If you can find such Images AVAILABLE in 'Input Documents', then incorporate them in the Media Block or Blocks and use their description for the the Media Block or Blocks. Alternatively, IF such images are NOT AVAILABLE in 'Input Documents', then USE YOUR IMAGINATION to create a Media Block or Blocks relevant to the text in the scenario and mention the type of Media (Image, Video, 360-Image, Audio) with description of its content and relevant overlay Tags for elaborating information and give directions to the course instructor of how to shoot and prepare these Media Blocks.)
    SimpleBranchingBlock (To select from a choice of choices (Branches) )
    Branch 1,2,3... (DIVISIBLE type containing path to other Branches) => with its TextBlock/s or None,MediaBlock/s or None, Branching Block (Simple Branching)
    Branch 1,2,3... (NON-DIVISIBLE type that are end of scenario branches not divisible further) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, Goal Block,  Debriefing, Reflection
    Branch 1,2,3... (NON-DIVISIBLE-MERGE type to link scenario branches when one story directly advances another branch's storyline) =>with its FeedbackAndFeedforwardBlock, TextBlock/s or None,MediaBlock/s or None, JumpBlock
    \nEnd of Overview structure\n

    Problems to overcome: 
    1. Produce a Media rich and diverse scenario by employing MediaBlock/s at various strategic places in the Scenario (specially Image type Media with overlayed hotspots), to add illustrativeness and elaborates content of the Text Blocks illustratively and visually presents the Choices in the Branching Blocks!, 
    2. 'timer' is only used for Text Blocks and Branching Blocks and the length of time is proportional to the content length in respective individual Text Blocks where timer is used.
        The decision time required in the Branching Blocks can be challenging or easy randomly, so base the length of the time according to the pertinent individual Branching Blocks.   

    SAMPLE EXAMPLE:::
{{
    "title": "(Insert a fitting Title Here)",
    "nodes": [
        {{
            "id": "StartBlock",
            "type": "StartBlock"
        }},
        {{
            "id": "B1",
            "type": "TextBlock",
            "title": "Learning_Objectives",
            "description": "1. (Insert Text Here); 2. (Insert Text Here) and so on"
        }},
        {{
            "id": "B2",
            "type": "TextBlock",
            "title": "Content_Areas",
            "description": "1. (Insert Text Here); 2. (Insert Text Here); 3. (Insert Text Here) and so on"
        }},
        {{
            "id": "B3",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "Bnhiefing of this Simulation Scenario",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "B4",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"The SBB below means SimpleBranchingBlock. The Bnh1, Bnh2 and so on are the branches.
        SBB_Bnh2 for example suggests it is the second branch from the SBB block."}},
        {{
            "id": "SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2": "(Insert Text Here) (DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh1_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh1_JB in this case is a story path that lead nowhere and brought the player back to the previous branching block SBB"}},
        {{
            "id": "SBB_Bnh1_JB",
            "type": "JumpBlock",
            "title": "Reevaluate Your Choices",
            "proceedToBlock": "SBB"
        }},
        {{
            "id": "SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_B3",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment":"SBB_Bnh2_SBB_Bnh3 for example suggests, if read and traced from backwards, it is the Third branch from the SBB block which
        in turn is from a Second branch that came from the very first SBB."}},
        {{
            "id": "SBB_Bnh2_SBB",
            "timer": "(Insert time in format hh:mm:ss)",
            "Purpose": "This block is where you !Divide the Simulation Game content into choices, that users can select and the corresponding divided branches leads to a consequence of the choice selected.",
            "type": "SimpleBranchingBlock",
            "title": "(Insert Text Here)",
            "branches": [
                {{
                    "port": "1",
                    "SBB_Bnh2_SBB_Bnh1": "(Insert Text Here) (NON-DIVISIBLE)"
                }},
                {{
                    "port": "2",
                    "SBB_Bnh2_SBB_Bnh2": "(Insert Text Here) (NON-DIVISIBLE-MERGE)"
                }},
                {{
                    "port": "3",
                    "SBB_Bnh2_SBB_Bnh3": "(Insert Text Here) (NON-DIVISIBLE)"
                }}
            ]
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh1_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_B2",
            "type": "MediaBlock",
            "title": "(Insert Text Here)",
            "mediaType": "Image (Preferred)/ 360-image/ Video/ Audio (Give one of these in your response)",
            "description": "(Insert Text Here)",
            "overlayTags": [
                "(Insert Text Here)"
            ]
        }},
        {{"_comment": "Jump blocks can be used for different reasons. Below SBB_Bnh2_SBB_Bnh2_JB in this case is a story path that lead the player to same outcome as another branch's goal block result of Bnh2_Bnh_Bnh3. Logically, it is possible that two paths taken by player can lead to a same outcome"}},
        {{
            "id": "SBB_Bnh2_SBB_Bnh2_JB",
            "type": "JumpBlock",
            "title": "(Insert Text Here)",
            "proceedToBlock": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B1",
            "type": "TextBlock",
            "title": "Feedback_And_Feedforward",
            "description": "Feedback=(Insert Text Here); Feedforward=(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_B2",
            "timer": "(Insert time in format hh:mm:ss)",
            "type": "TextBlock",
            "title": "(Insert Text Here)",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_GB",
            "type": "GoalBlock",
            "title": "(Insert Text Here)",
            "score": "Insert Integer Number Here. Give smaller score then the relevant Correct Choice Bnhanch score"
        }},
        {{
            "id": "BSBB_Bnh2_SBB_Bnh3_DB",
            "type": "TextBlock",
            "title": "Debriefing",
            "description": "(Insert Text Here)"
        }},
        {{
            "id": "SBB_Bnh2_SBB_Bnh3_RF",
            "type": "TextBlock",
            "title": "Reflection",
            "description": "(Insert Text Here)"
        }}
    ],                       
    "edges": [
        {{
            "source": "StartBlock",
            "target": "B1"
        }},
        {{
            "source": "B1",
            "target": "B2"
        }},
        {{
            "source": "B2",
            "target": "B3"
        }},
        {{
            "source": "B3",
            "target": "B4"
        }},
        {{
            "source": "B4",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh1_B1",
            "sourceport": "1"
        }},
        {{
            "source": "SBB_Bnh1_B1",
            "target": "SBB_Bnh1_B2"
        }},
        {{
            "source": "SBB_Bnh1_B2",
            "target": "SBB_Bnh1_JB"
        }},
        {{
            "source": "SBB_Bnh1_JB",
            "target": "SBB"
        }},
        {{
            "source": "SBB",
            "target": "SBB_Bnh2_B1",
            "sourceport": "2"
        }},
        {{
            "source": "SBB_Bnh2_B1",
            "target": "SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_B2",
            "target": "SBB_Bnh2_B3"
        }},
        {{
            "source": "SBB_Bnh2_B3",
            "target": "SBB_Bnh2_SBB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh1_B1",
            "sourceport":"1"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_B1",
            "target": "SBB_Bnh2_SBB_Bnh1_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_GB",
            "target": "SBB_Bnh2_SBB_Bnh1_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh1_DB",
            "target": "SBB_Bnh2_SBB_Bnh1_RF"
        }}
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh2_B1",
            "sourceport":"2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B1",
            "target": "SBB_Bnh2_SBB_Bnh2_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_B2",
            "target": "SBB_Bnh2_SBB_Bnh2_JB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh2_JB",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB",
            "target": "SBB_Bnh2_SBB_Bnh3_B1",
            "sourceport":"3"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B1",
            "target": "SBB_Bnh2_SBB_Bnh3_B2"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_B2",
            "target": "SBB_Bnh2_SBB_Bnh3_GB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_GB",
            "target": "SBB_Bnh2_SBB_Bnh3_DB"
        }},
        {{
            "source": "SBB_Bnh2_SBB_Bnh3_DB",
            "target": "SBB_Bnh2_SBB_Bnh3_RF"
        }}
    ]
}}
    SAMPLE EXAMPLE END

    !!!ATTENTION!!!
    Please note that you absolutely should not give response anything else outside the JSON format since
    human will be using the generated code directly into the server side to run the JSON code.
    Moreover, it is absolutley mandatory and necessary for you to generate a complete JSON response such that the JSON generated from you must enclose all the parenthesis at the end of your response
    and all it's parameters are also closed in the required syntax rules of JSON and all the blocks be included in it since we want our JSON
    to be compilable. 
    You Prefer to make simulation such that a choice may lead to a consequnece that may lead to more choice or choices that may lead to more consequences, evetually reaching the end of the scenario.
    Give concise, relevant, clear, and descriptive instructions as you are an educational provider that has expertise 
    in molding asked information into the said block structure to teach and instruct students.     

    NEGATIVE PROMPT: Responding outside the JSON format.   

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.
    ]

    
    Chatbot:"""
)
### Simulation Prompts End

prompt_LO_CA = PromptTemplate(
    input_variables=["input_documents","human_input"],
    template="""
    Based on the information provided in 'Human Input' and 'Input Documents', you are going to generate 
    Learning Objectives and Content Areas in a JSON format. Make sure the both Learning Objectives and Content Areas
    are specifically relevant to the query of 'Human Input'. 
    Lets suppose the 'Human Input' asks for a course to be created for Driving a Car. 
    And the 'Input Documents' has information for both Driving a Car and Repairing a Car. 
    Then, you should only give the Learning Objectives and Content Areas about Driving a Car only
    since the 'Human Input' asked you only about this topic.

    Do not give any Learning Objectives or Content Areas based on information
    not present in the 'Input Documents'. You have to just strictly keep the Learning Objectives and Content Areas
    limited and specific to the information asked by 'Human Input' AND present in the 'Input Documents'; and nothing outside it.
    ***Stick strictly to the information given in the 'Input Documents' provided to you.
    The 'Human Input' decides what information to collect from the 'Input Documents' to create Learning Objectives
    and Content Areas.***

    *DIRE WARNING: The number of points of Learning Objectives and Content Areas can be different.
    The Example below is only given for context of format and absolutely NOT for the fact that you 
    generate same number of points as given in the Example for the Learning Objectives and Content Areas. 
    Learning Objectives and Content Areas can have only 1 point or more points, 
    all depends on the amount of information present in the 'Input Documents'
    and the query pertaining to it by the human in the 'Human Input'.*
    
    \nExample\n
    {{
    "LearningObjectives": [
        "1. Recognize the Signs and Symptoms of a Heart Attack: Learners will be able to identify both typical and atypical signs of a heart attack, understanding that symptoms can vary widely among individuals.\n2. Emergency Response Procedures: Learners will understand the steps to take in both scenarios where the patient is unconscious and conscious, including the use of DRSABCD (Danger, Response, Send for help, Airway, Breathing, CPR, Defibrillation)."
    ],
    "ContentAreas": [
        "1. Introduction to Heart Attacks: Overview of what constitutes a heart attack, including the physiological underpinnings and the importance of quick response.\n2. Identifying Symptoms: Detailed review of both common and less common symptoms of heart attacks, with emphasis on variations by gender, age, and pre-existing conditions.\n3. First Aid Steps: Step-by-step guide for responding to a heart attack in various situations (unconscious vs. conscious patients)."
    ]
    }}
    \nExample End\n

    'Human Input': {human_input}
    'Input Documents': {input_documents}
    Chatbot:"""
)

prompt_LO_CA_GEMINI = PromptTemplate(
    input_variables=["input_documents","human_input"],
    template="""
    Based on the information provided in 'Human Input' and 'Input Documents', you are going to generate 
    Learning Objectives and Content Areas in a JSON format. Make sure the both Learning Objectives and Content Areas
    are specifically relevant to the query of 'Human Input'. 
    
    \nYour Example Response Format\n
    {{
    "LearningObjectives": [
        "Insert Learning Objective or Objectives here"
    ],
    "ContentAreas": [
        "Insert Content Area or Areas in here"
    ]
    }}
    \nExample End\n

    'Human Input': {human_input}
    'Input Documents': {input_documents}

    DO NOT START YOUR RESPONSE WITH ```json and END WITH ``` 
    Just start the JSON response directly.
    """
)

def PRODUCE_LEARNING_OBJ_COURSE(query, docsearch, llm, model_type):
    print("PRODUCE_LEARNING_OBJ_COURSE Initiated!")
    docs = docsearch.similarity_search(query, k=3)
    docs_main = " ".join([d.page_content for d in docs])
    if model_type == "gemini": 
        print("Now processing prompt_LO_CA_GEMINI")
        chain = LLMChain(prompt=prompt_LO_CA_GEMINI, llm=llm)
    else:
        chain = LLMChain(prompt=prompt_LO_CA, llm=llm)
    return chain, docs_main, query

def RE_SIMILARITY_SEARCH(query, docsearch, output_path, model_type, summarize_images):
    print("RE_SIMILARITY_SEARCH Initiated!")
    docs = docsearch.similarity_search(query, k=3)
    print("docs from RE_SIMILARITY_SEARCH",docs)
    if summarize_images == "on":
        print(f"Tells me to summarize images, {summarize_images}")
        PageNumberList = []
        for relevant_doc in docs:
            relevant_doc = relevant_doc.page_content
            print(relevant_doc)

            pattern_this_pptx = r"'SlideNumber': (\d+), 'FileName': '(.+?)'"        # f'SlideNumber:{slide_number} of FileName:{filename_without_extension}-ImageNumber {image_number}'
            # Find all matches for "[This Page is PageNumber:]"
            matches_this_pptx = re.findall(pattern_this_pptx, relevant_doc)

            pattern_this_end_pptx = r"End of SlideNumber:(\d+) with Filename:(.+?) ----"        # f'SlideNumber:{slide_number} of FileName:{filename_without_extension}-ImageNumber {image_number}'
            # Find all matches for "[This Page is PageNumber:]"
            matches_this_end_pptx = re.findall(pattern_this_end_pptx, relevant_doc)


            pattern_this_doc = r'----media/ImageNumber:(\d+) PageNumber:Null of FileName:(.+)----'
            matches_this_doc = re.findall(pattern_this_doc, relevant_doc)
            
            pattern_end = r'End of PageNumber:(\d+) of file name:(.+)\n'
            pattern_this_page = r'The Content of PageNumber:(\d+) of file name:(.+) is:\n'
            # Find all matches for "End of PageNumber:"
            matches_end = re.findall(pattern_end, relevant_doc)

            # Find all matches for "[This Page is PageNumber:]"
            matches_this_page = re.findall(pattern_this_page, relevant_doc)

            # Combine the matches
            for num in matches_this_page + matches_end + matches_this_doc + matches_this_pptx + matches_this_end_pptx:
                PageNumberList.append(num)

            PageNumberList = list(set(PageNumberList))
            print("PageNumberList",PageNumberList)

        image_elements = []
        image_summaries = []

        def encode_image(image_path):
            basename = os.path.basename(image_path)
            with Image.open(image_path) as img:
                width, height = img.size
                print(f"{basename} size is {width},{height}")
                if width*height > 262144:
                    # Resize the image
                    img = img.resize((512, 512))
                    # Save the resized image to a temporary file
                    basenama = os.path.basename(image_path)
                    extensiona = basenama.rsplit('.', 1)[1].lower()

                    temp_patha = image_path + f"_temp_img.{extensiona}"
                    img.save(temp_patha)
                
                    # Encode the resized image
                    with open(temp_patha, "rb") as f:
                        encoded_image = base64.b64encode(f.read()).decode('utf-8')

                    # Remove the temporary file
                    os.remove(temp_patha)
                else:
                    print(f"{basename} is less than 262144 having {width}, {height}")
                    with open(image_path, "rb") as f:
                        encoded_image = base64.b64encode(f.read()).decode('utf-8')

                return encoded_image



        def summarize_image(encoded_image, basename):
            prompt = [
                SystemMessage(content="You are a bot that is good at analyzing images."),
                HumanMessage(content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image. Tell what FileName, PageNumber/SlideNumber and ImageNumber of this image is by seeing this information: {basename}. Your output should look like this: 'This image that belongs to FileName: ..., PageNumber: ..., ImageNumber: .... In this Image ...' or in case of SlideNumber available 'This image that belongs to FileName: ..., SlideNumber: ..., ImageNumber: .... In this Image ...' !!!WARNING: Exact, absolutely Unchanged File name of the image must be mentioned as found in {basename}. File name may contain special characters such as hyphens (-), underscores (_), semicolons (;), spaces, and others, so this should be kept in mind!!!"
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{encoded_image}"
                        },
                    },
                ])
            ]

            prompt_gemini = HumanMessage(
                content=[
                    {
                        "type": "text",
                        "text": f"Describe the contents of this image. Tell what FileName, PageNumber/SlideNumber and ImageNumber of this image is by seeing this information: {basename}. Your output should look like this: 'This image that belongs to FileName: ..., PageNumber: ..., ImageNumber: .... In this Image ...' or in case of SlideNumber available 'This image that belongs to FileName: ..., SlideNumber: ..., ImageNumber: .... In this Image ...' !!!WARNING: Exact, absolutely Unchanged File name of the image must be mentioned as found in {basename}. File name may contain special characters such as hyphens (-), underscores (_), semicolons (;), spaces, and others, so this should be kept in mind!!!",
                    },  # You can optionally provide text parts
                    {"type": "image_url", "image_url": f"data:image/jpeg;base64,{encoded_image}"},
                ]
            )

            if model_type == 'gemini':
                print("Gemini summarizing images NOW")
                response = ChatGoogleGenerativeAI(model="gemini-pro-vision",temperature=0,max_output_tokens=200).invoke([prompt_gemini])
                return response.content
            else:
                response = ChatOpenAI(model="gpt-4o", max_tokens=200, temperature=0).invoke(prompt)
                return response.content

        for root, dirs, files in os.walk(output_path):
            for i in files:
                if i.endswith(('.png', '.jpg', '.jpeg')):
                    for page_number, file in PageNumberList:
                        if f"FileName {file} PageNumber {page_number}" in i:
                            image_path = os.path.join(root, i)
                            basename = os.path.basename(image_path)
                            print(os.path.basename(image_path))
                            encoded_image = encode_image(image_path)
                            image_elements.append(encoded_image)
                            summary = summarize_image(encoded_image,basename)
                            image_summaries.append(summary)
                        elif f"FileName {file} PageNumber Null ImageNumber {page_number}" in i:
                            image_path = os.path.join(root, i)
                            basename = os.path.basename(image_path)
                            print(os.path.basename(image_path))
                            encoded_image = encode_image(image_path)
                            image_elements.append(encoded_image)
                            summary = summarize_image(encoded_image,basename)
                            image_summaries.append(summary)
                        elif f"FileName {file} SlideNumber {page_number}" in i:
                            image_path = os.path.join(root, i)
                            basename = os.path.basename(image_path)
                            print(os.path.basename(image_path))
                            encoded_image = encode_image(image_path)
                            image_elements.append(encoded_image)
                            summary = summarize_image(encoded_image,basename)
                            image_summaries.append(summary)

        print("image_summaries::",image_summaries)

        image_summaries_string = "\n".join(image_summaries) #convert list to string to add in the langchain Document data type
        docs.append(Document(page_content=f"Useful Image/s for all the above content::\n{image_summaries_string}"))

    return docs



def TALK_WITH_RAG(scenario, content_areas, learning_obj, query, docs_main, llm, model_type, model_name):
    print("TALK_WITH_RAG Initiated!")
    # if we are getting docs_main already from the process_data flask route then comment, else
    # UNcomment if you want more similarity_searching based on Learning obj and content areas!
    # docs = docsearch.similarity_search(query, k=3)
    # docs_main = " ".join([d.page_content for d in docs])
    responses = ''
    def is_json_parseable(json_string):
        try:
            json_object = json.loads(json_string)
        except ValueError as e:
            return False, str(e)
        return True, json_object
         
    if scenario == "linear":
        print("SCENARIO ====prompt_linear",scenario)
        # if model_type == 'gemini':
        #     chain = LLMChain(prompt=prompt_linear_gemini, llm=llm)
        # else:
        #     chain = LLMChain(prompt=prompt_linear, llm=llm)
        
        chain = LLMChain(prompt=prompt_linear, llm=llm)
        response = chain({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})

        is_valid, result = is_json_parseable(response['text'])
        
        if is_valid == False:
            txt = response['text']
            print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
            ### REGEX to remove last incomplete id block ###
            # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
            
            y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
            matches = re.findall(y, txt, re.DOTALL)
            print(len(matches))
            print(matches[-1])

            if matches:
                last_match = matches[-1]
                idx = txt.rfind(last_match)  # find the index of the last occurrence
                modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                print("Original text:")
                print(txt)

                print("\nModified text:")
                print(modified_txt)
            else:
                print("No matches found.")
                ### ###
            responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

            chain_retry = LLMChain(prompt=prompt_linear_retry,llm=llm)
            response_retry = chain_retry({"incomplete_response": modified_txt})
            print("response contd... is:\n",response_retry['text'])

            responses = modified_txt + response_retry['text']
            print("Responses list is:\n",responses)
            is_valid_retry, result = is_json_parseable(responses)
            if is_valid_retry == False:
                print("The retry is also not parseable!", responses)
                max_attempts = 3  # Maximum number of attempts
                attempts = 1
                while attempts < max_attempts:
                    chain_simplify = LLMChain(prompt=prompt_linear_simplify,llm=llm)
                    response_retry_simplify = chain_simplify({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                    is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                    if is_valid_retry_simplify == True:
                        response['text'] = response_retry_simplify['text']
                        print("Result successfull for simplified response:",response['text'])
                        break
                    else:
                        attempts += 1
                        print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
            else:
                response['text'] = responses
                print("Retry success", response['text'])
        else:
            ("Parseable JSON", response['text'])

    elif scenario == "branched":
        print("SCENARIO ====branched",scenario)
        
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
        else:
            llm_setup = ChatOpenAI(model=model_name, temperature=0)

        # summarized first, then response
        chain1 = LLMChain(prompt=prompt_branched_setup,llm=llm_setup)
        response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        print("Response 1 is::",response1['text'])
    
        chain = LLMChain(prompt=prompt_branched,llm=llm)  
        response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})

        is_valid, result = is_json_parseable(response['text'])
        
        if is_valid == False:
            txt = response['text']
            print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
            ### REGEX to remove last incomplete id block ###
            # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
            
            y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
            matches = re.findall(y, txt, re.DOTALL)
            print(len(matches))
            print(matches[-1])

            if matches:
                last_match = matches[-1]
                idx = txt.rfind(last_match)  # find the index of the last occurrence
                modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                print("Original text:")
                print(txt)

                print("\nModified text:")
                print(modified_txt)
            else:
                print("No matches found.")
                ### ###
            responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

            chain_retry = LLMChain(prompt=prompt_branched_retry,llm=llm)
            response_retry = chain_retry({"incomplete_response": modified_txt,"micro_subtopics":response1['text']})
            print("response contd... is:\n",response_retry['text'])

            responses = modified_txt + response_retry['text']
            print("Responses list is:\n",responses)
            is_valid_retry, result = is_json_parseable(responses)
            if is_valid_retry == False:
                print("The retry is also not parseable!", responses)
                max_attempts = 3  # Maximum number of attempts
                attempts = 1
                while attempts < max_attempts:
                    chain_simplify = LLMChain(prompt=prompt_branched_simplify,llm=llm)
                    response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                    is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                    if is_valid_retry_simplify == True:
                        response['text'] = response_retry_simplify['text']
                        print("Result successfull for simplified response:",response['text'])
                        break
                    else:
                        attempts += 1
                        print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
            else:
                response['text'] = responses
                print("Retry success", response['text'])
        else:
            ("Parseable JSON", response['text'])

    elif scenario == "simulation":
        print("SCENARIO ====prompt_simulation_pedagogy",scenario)
        # summarized first, then response
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0.3)
            chain = LLMChain(prompt=prompt_simulation_pedagogy_gemini,llm=llm)
            print("prompt_simulation_pedagogy_gemini selected!")
        else:
            llm_setup = ChatOpenAI(model=model_name, temperature=0.3)
            chain = LLMChain(prompt=prompt_simulation_pedagogy_gemini,llm=llm)

        chain1 = LLMChain(prompt=prompt_simulation_pedagogy_setup,llm=llm_setup)
        response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        print("Response 1 is::",response1['text'])

        response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        
        is_valid, result = is_json_parseable(response['text'])
        
        if is_valid == False:
            txt = response['text']
            print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
            ### REGEX to remove last incomplete id block ###
            # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
            
            y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
            matches = re.findall(y, txt, re.DOTALL)
            print(len(matches))
            print(matches[-1])

            if matches:
                last_match = matches[-1]
                idx = txt.rfind(last_match)  # find the index of the last occurrence
                modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                print("Original text:")
                print(txt)

                print("\nModified text:")
                print(modified_txt)
            else:
                print("No matches found.")
                ### ###
            responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

            chain_retry = LLMChain(prompt=prompt_simulation_pedagogy_retry_gemini,llm=llm)
            response_retry = chain_retry({"incomplete_response": modified_txt,"simulation_story":response1['text']})
            print("response contd... is:\n",response_retry['text'])

            responses = modified_txt + response_retry['text']
            print("Responses list is:\n",responses)
            is_valid_retry, result = is_json_parseable(responses)
            if is_valid_retry == False:
                print("The retry is also not parseable!", responses)
                max_attempts = 3  # Maximum number of attempts
                attempts = 1
                while attempts < max_attempts:
                    chain_simplify = LLMChain(prompt=prompt_simulation_pedagogy_gemini_simplify,llm=llm)
                    response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                    is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                    if is_valid_retry_simplify == True:
                        response['text'] = response_retry_simplify['text']
                        print("Result successfull for simplified response:",response['text'])
                        break
                    else:
                        attempts += 1
                        print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
            else:
                response['text'] = responses
                print("Retry success", response['text'])
        else:
            ("Parseable JSON", response['text'])
            
    elif scenario == "gamified":
        print("SCENARIO ====prompt_gamified",scenario)
        if model_type == 'gemini':
            llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
        else:
            llm_setup = ChatOpenAI(model=model_name, temperature=0)

        chain1 = LLMChain(prompt=prompt_gamified_setup,llm=llm_setup)
        response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        print("Response 1 is::",response1['text'])

        # chain2 = LLMChain(prompt=prompt_gamified_simple,llm=llm_setup)
        # response2 = chain2({"response_of_bot_simple": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        # print("Response 2 is::",response2['text'])

        chain = LLMChain(prompt=prompt_gamified_json,llm=llm)
        response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
        
        is_valid, result = is_json_parseable(response['text'])
        
        if is_valid == False:
            txt = response['text']
            print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
            ### REGEX to remove last incomplete id block ###
            # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
            y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
            matches = re.findall(y, txt, re.DOTALL)
            print(len(matches))
            print(matches[-1])

            if matches:
                last_match = matches[-1]
                idx = txt.rfind(last_match)  # find the index of the last occurrence
                modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                print("Original text:")
                print(txt)

                print("\nModified text:")
                print(modified_txt)
            else:
                print("No matches found.")
                ### ###
            responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

            chain_retry = LLMChain(prompt=prompt_gamified_pedagogy_retry_gemini,llm=llm)
            response_retry = chain_retry({"incomplete_response": modified_txt,"exit_game_story":response1['text']})
            print("response contd... is:\n",response_retry['text'])

            responses = modified_txt + response_retry['text']
            print("Responses list is:\n",responses)
            is_valid_retry, result = is_json_parseable(responses)
            if is_valid_retry == False:
                print("The retry is also not parseable!", responses)
                max_attempts = 3  # Maximum number of attempts
                attempts = 1
                while attempts < max_attempts:
                    chain_simplify = LLMChain(prompt=prompt_gamify_pedagogy_gemini_simplify,llm=llm)
                    response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                    is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                    if is_valid_retry_simplify == True:
                        response['text'] = response_retry_simplify['text']
                        print("Result successfull for simplified response:",response['text'])
                        break
                    else:
                        attempts += 1
                        print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
            else:
                response['text'] = responses
                print("Retry success", response['text'])
        else:
            ("Parseable JSON", response['text'])

    elif scenario == "auto":
        print("SCENARIO ====PROMPT",scenario)
        # chain = prompt | llm | {f"{llm_memory}": RunnablePassthrough()}
        

        ### SEMANTIC ROUTES LOGIC ###
        if model_type == 'gemini':
            llm_auto = ChatGoogleGenerativeAI(model=model_name,temperature=0.4, max_output_tokens=32)
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        else:
            embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
            llm_auto = ChatOpenAI(model=model_name, temperature=0.4, max_tokens=32)
        
        llm_auto_chain = LLMChain(prompt=promptSelector, llm=llm_auto)
        selected = llm_auto_chain.run({"input_documents": docs_main, "human_input": query})

        print("Semantic Scenario Selected of NAME",selected)

        gamified_route = ['gamified', 'gamified scenario','bot: gamified scenario']
        simulation_route = ['simulation', 'simulation scenario', 'bot: simulation scenario']
        linear_route = ['linear', 'linear scenario', 'bot: linear scenario']
        branched_route = ['branched', 'branched scenario', 'bot: branched scenario']

        gamified_route_embeddings = embeddings.embed_documents(gamified_route)
        simulation_route_embeddings = embeddings.embed_documents(simulation_route)
        linear_route_embeddings = embeddings.embed_documents(linear_route)
        branched_route_embeddings =  embeddings.embed_documents(branched_route)

        query_embedding = embeddings.embed_query(selected)

        gamified_similarity = cosine_similarity([query_embedding],gamified_route_embeddings)[0]
        simulation_similarity = cosine_similarity([query_embedding],simulation_route_embeddings)[0]
        linear_similarity = cosine_similarity([query_embedding], linear_route_embeddings)[0]
        branched_similarity = cosine_similarity([query_embedding], branched_route_embeddings)[0]

        max_similarity = max(max(gamified_similarity), max(simulation_similarity), max(linear_similarity), max(branched_similarity))

        ############################
        
        if max_similarity == max(gamified_similarity):
            print("Gamified Auto Selected")
            if model_type == 'gemini':
                llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
            else:
                llm_setup = ChatOpenAI(model=model_name, temperature=0)

            chain1 = LLMChain(prompt=prompt_gamified_setup,llm=llm_setup)
            response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            print("Response 1 is::",response1['text'])

            # chain2 = LLMChain(prompt=prompt_gamified_simple,llm=llm_setup)
            # response2 = chain2({"response_of_bot_simple": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            # print("Response 2 is::",response2['text'])

            chain = LLMChain(prompt=prompt_gamified_json,llm=llm)
            response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            
            is_valid, result = is_json_parseable(response['text'])
            
            if is_valid == False:
                txt = response['text']
                print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
                ### REGEX to remove last incomplete id block ###
                # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
                y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
                matches = re.findall(y, txt, re.DOTALL)
                print(len(matches))
                print(matches[-1])

                if matches:
                    last_match = matches[-1]
                    idx = txt.rfind(last_match)  # find the index of the last occurrence
                    modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                    print("Original text:")
                    print(txt)

                    print("\nModified text:")
                    print(modified_txt)
                else:
                    print("No matches found.")
                    ### ###
                responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

                chain_retry = LLMChain(prompt=prompt_gamified_pedagogy_retry_gemini,llm=llm)
                response_retry = chain_retry({"incomplete_response": modified_txt,"exit_game_story":response1['text']})
                print("response contd... is:\n",response_retry['text'])

                responses = modified_txt + response_retry['text']
                print("Responses list is:\n",responses)
                is_valid_retry, result = is_json_parseable(responses)
                if is_valid_retry == False:
                    print("The retry is also not parseable!", responses)
                    max_attempts = 3  # Maximum number of attempts
                    attempts = 1
                    while attempts < max_attempts:
                        chain_simplify = LLMChain(prompt=prompt_gamify_pedagogy_gemini_simplify,llm=llm)
                        response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                        is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                        if is_valid_retry_simplify == True:
                            response['text'] = response_retry_simplify['text']
                            print("Result successfull for simplified response:",response['text'])
                            break
                        else:
                            attempts += 1
                            print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
                else:
                    response['text'] = responses
                    print("Retry success", response['text'])
            else:
                ("Parseable JSON", response['text'])

        elif max_similarity == max(linear_similarity):
            print("Linear Auto Selected")
            chain = LLMChain(prompt=prompt_linear, llm=llm)
            response = chain({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})

            is_valid, result = is_json_parseable(response['text'])
            
            if is_valid == False:
                txt = response['text']
                print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
                ### REGEX to remove last incomplete id block ###
                # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
                
                y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
                matches = re.findall(y, txt, re.DOTALL)
                print(len(matches))
                print(matches[-1])

                if matches:
                    last_match = matches[-1]
                    idx = txt.rfind(last_match)  # find the index of the last occurrence
                    modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                    print("Original text:")
                    print(txt)

                    print("\nModified text:")
                    print(modified_txt)
                else:
                    print("No matches found.")
                    ### ###
                responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

                chain_retry = LLMChain(prompt=prompt_linear_retry,llm=llm)
                response_retry = chain_retry({"incomplete_response": modified_txt})
                print("response contd... is:\n",response_retry['text'])

                responses = modified_txt + response_retry['text']
                print("Responses list is:\n",responses)
                is_valid_retry, result = is_json_parseable(responses)
                if is_valid_retry == False:
                    print("The retry is also not parseable!", responses)
                    max_attempts = 3  # Maximum number of attempts
                    attempts = 1
                    while attempts < max_attempts:
                        chain_simplify = LLMChain(prompt=prompt_linear_simplify,llm=llm)
                        response_retry_simplify = chain_simplify({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                        is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                        if is_valid_retry_simplify == True:
                            response['text'] = response_retry_simplify['text']
                            print("Result successfull for simplified response:",response['text'])
                            break
                        else:
                            attempts += 1
                            print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
                else:
                    response['text'] = responses
                    print("Retry success", response['text'])
            else:
                ("Parseable JSON", response['text'])

        elif max_similarity == max(simulation_similarity):
            print("Simulation Auto Selected")
            if model_type == 'gemini':
                llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0.3)
                chain = LLMChain(prompt=prompt_simulation_pedagogy_gemini,llm=llm)
                print("prompt_simulation_pedagogy_gemini selected!")
            else:
                llm_setup = ChatOpenAI(model=model_name, temperature=0.3)
                chain = LLMChain(prompt=prompt_simulation_pedagogy_gemini,llm=llm)

            chain1 = LLMChain(prompt=prompt_simulation_pedagogy_setup,llm=llm_setup)
            response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            print("Response 1 is::",response1['text'])

            response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            
            is_valid, result = is_json_parseable(response['text'])
            
            if is_valid == False:
                txt = response['text']
                print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
                ### REGEX to remove last incomplete id block ###
                # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
                
                y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
                matches = re.findall(y, txt, re.DOTALL)
                print(len(matches))
                print(matches[-1])

                if matches:
                    last_match = matches[-1]
                    idx = txt.rfind(last_match)  # find the index of the last occurrence
                    modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                    print("Original text:")
                    print(txt)

                    print("\nModified text:")
                    print(modified_txt)
                else:
                    print("No matches found.")
                    ### ###
                responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

                chain_retry = LLMChain(prompt=prompt_simulation_pedagogy_retry_gemini,llm=llm)
                response_retry = chain_retry({"incomplete_response": modified_txt,"simulation_story":response1['text']})
                print("response contd... is:\n",response_retry['text'])

                responses = modified_txt + response_retry['text']
                print("Responses list is:\n",responses)
                is_valid_retry, result = is_json_parseable(responses)
                if is_valid_retry == False:
                    print("The retry is also not parseable!", responses)
                    max_attempts = 3  # Maximum number of attempts
                    attempts = 1
                    while attempts < max_attempts:
                        chain_simplify = LLMChain(prompt=prompt_simulation_pedagogy_gemini_simplify,llm=llm)
                        response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                        is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                        if is_valid_retry_simplify == True:
                            response['text'] = response_retry_simplify['text']
                            print("Result successfull for simplified response:",response['text'])
                            break
                        else:
                            attempts += 1
                            print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
                else:
                    response['text'] = responses
                    print("Retry success", response['text'])
            else:
                ("Parseable JSON", response['text'])

        elif max_similarity == max(branched_similarity):
            print("Branched Auto Selected")
            if model_type == 'gemini':
                llm_setup = ChatGoogleGenerativeAI(model=model_name,temperature=0)
            else:
                llm_setup = ChatOpenAI(model=model_name, temperature=0)

            # summarized first, then response
            chain1 = LLMChain(prompt=prompt_branched_setup,llm=llm_setup)
            response1 = chain1({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
            print("Response 1 is::",response1['text'])
        
            chain = LLMChain(prompt=prompt_branched,llm=llm)  
            response = chain({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})

            is_valid, result = is_json_parseable(response['text'])
            
            if is_valid == False:
                txt = response['text']
                print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
                ### REGEX to remove last incomplete id block ###
                # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
                
                y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
                matches = re.findall(y, txt, re.DOTALL)
                print(len(matches))
                print(matches[-1])

                if matches:
                    last_match = matches[-1]
                    idx = txt.rfind(last_match)  # find the index of the last occurrence
                    modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                    print("Original text:")
                    print(txt)

                    print("\nModified text:")
                    print(modified_txt)
                else:
                    print("No matches found.")
                    ### ###
                responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

                chain_retry = LLMChain(prompt=prompt_branched_retry,llm=llm)
                response_retry = chain_retry({"incomplete_response": modified_txt,"micro_subtopics":response1['text']})
                print("response contd... is:\n",response_retry['text'])

                responses = modified_txt + response_retry['text']
                print("Responses list is:\n",responses)
                is_valid_retry, result = is_json_parseable(responses)
                if is_valid_retry == False:
                    print("The retry is also not parseable!", responses)
                    max_attempts = 3  # Maximum number of attempts
                    attempts = 1
                    while attempts < max_attempts:
                        chain_simplify = LLMChain(prompt=prompt_branched_simplify,llm=llm)
                        response_retry_simplify = chain_simplify({"response_of_bot": response1['text'],"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                        is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                        if is_valid_retry_simplify == True:
                            response['text'] = response_retry_simplify['text']
                            print("Result successfull for simplified response:",response['text'])
                            break
                        else:
                            attempts += 1
                            print(f"Attempt {attempts}: Failed to parse JSON. Error:\n {response_retry_simplify['text']}")
                else:
                    response['text'] = responses
                    print("Retry success", response['text'])
            else:
                ("Parseable JSON", response['text'])        

        else:
            print("AUTO SELECTION FAILED, Selecting Default Scenario of LINEAR SCENARIO")

            chain = LLMChain(prompt=prompt_linear, llm=llm)
            response = chain({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})

            is_valid, result = is_json_parseable(response['text'])
            
            if is_valid == False:
                txt = response['text']
                print("CHAIN_RETRY BEGINS for the failed response:\n", txt)
                ### REGEX to remove last incomplete id block ###
                # y = r'\s*{\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z)' # this is previous working iteration for the below regex
                
                y = r'\s*},\s*.*?({\s*"id":\s*"[^"]+".*?(?=\n\s*{|\Z))'
                matches = re.findall(y, txt, re.DOTALL)
                print(len(matches))
                print(matches[-1])

                if matches:
                    last_match = matches[-1]
                    idx = txt.rfind(last_match)  # find the index of the last occurrence
                    modified_txt = txt[:idx]  # remove the last occurrence and everything after it

                    print("Original text:")
                    print(txt)

                    print("\nModified text:")
                    print(modified_txt)
                else:
                    print("No matches found.")
                    ### ###
                responses = modified_txt + "[CONTINUE_EXACTLY_FROM_HERE]"

                chain_retry = LLMChain(prompt=prompt_linear_retry,llm=llm)
                response_retry = chain_retry({"incomplete_response": modified_txt})
                print("response contd... is:\n",response_retry['text'])

                responses = modified_txt + response_retry['text']
                print("Responses list is:\n",responses)
                is_valid_retry, result = is_json_parseable(responses)
                if is_valid_retry == False:
                    print("The retry is also not parseable!", responses)
                    max_attempts = 3  # Maximum number of attempts
                    attempts = 1
                    while attempts < max_attempts:
                        chain_simplify = LLMChain(prompt=prompt_linear_simplify,llm=llm)
                        response_retry_simplify = chain_simplify({"input_documents": docs_main,"human_input": query,"content_areas": content_areas,"learning_obj": learning_obj})
                        is_valid_retry_simplify, result = is_json_parseable(response_retry_simplify['text'])
                        if is_valid_retry_simplify == True:
                            response['text'] = response_retry_simplify['text']
                            print("Result successfull for simplified response:",response['text'])
                            break
                        else:
                            attempts += 1
                            print(f"Attempt {attempts}: Failed to parse JSON. Error: {response_retry_simplify['text']}")
                else:
                    response['text'] = responses
                    print("Retry success", response['text'])
            else:
                ("Parseable JSON", response['text'])
    
    print("The output is as follows::\n",response['text'])
    return response['text']

def ANSWER_IMG(response_text, llm,relevant_doc):
    # prompt_template_img =PromptTemplate( 
    # input_variables=["response_text","context"],
    # template="""
    # Provided the context, look at the Images that are mentioned in the 'response_text': {response_text}. Provide a brief summary of those 
    # images stored in the 'context': {context}.
    # Format of Reply (The number of Images and their description may vary, depends on what is instructed in the
    # 'response_text'. If only one image is mentioned in the 'response_text', then you should include Image1 only. If there are 2 or more images then your reply should
    # also have same images as mentioned in the 'response_text'!):
    # {{"Image1": "file_name_..._page_..._image_...",
    # "Description1": "...",
    # "Image2": "file_name_..._page_..._image_...",
    # "Description2": "..."
    # and so on
    # }}
    # Warning: Include the complete schema of name defined. The complete schema of name includes
    # "file_name_..._page_..._image_..."
    # Take great care for the underscores. They are to be used exactly as defined. Also take 
    # extreme caution at the file_name since the file might be having its own - and _ which is not to be
    # tampered with in any way and should remain exactly the same!
    # [WARNING: The ... presents page number as int and image number as int. But, for the file_name_ it represents
    # as the file name itself which may have its own dashes or underscores or brackets. Whatever the file name
    # you found in the 'context', make sure you use the same name. ]
    # Answer():
    # """
    # )



    # chain = LLMChain(prompt=prompt_template_img,llm=llm)
    # img_response = chain.run({"response_text": response_text, "context": relevant_doc})
    # print("img_response is::",img_response)
###

    class image_loc(BaseModel):
        FileName: str = Field(description="Exact, absolutely Unchanged File name of the image as mentioned in the 'Context'. File name may contain special characters such as hyphens (-), underscores (_), semicolons (;), spaces, and others.")
        PageNumber: Optional[str] = Field(description="If available, write page number of the image. 'Null' if not available. !!!DO NOT USE PageNumber if SlideNumber is available.!!!")
        SlideNumber: Optional[str] = Field(description="If available, slide number of the image.")
        ImageNumber: int = Field(description="image number of the image")
        Description: str = Field(description="Description detail of the image")

    class image(BaseModel):
        Image: List[image_loc] = Field(description="image_loc")

    parser = JsonOutputParser(pydantic_object=image)

    prompt = PromptTemplate(
    template="""
    Search for those image or images only, whose descriptions in the Media Blocks of the 'Response Text' matches
    with the descriptions in the 'Context' data. Output only those image's or images' description from the 
    'Context' data. 
    \n{format_instructions}\n'Response Text': {response_text}\n'Context': {context}""",
    input_variables=["response_text","context"],
    partial_variables={"format_instructions": parser.get_format_instructions()},
    )

    chain = prompt | llm | parser

    img_response = chain.invoke({"response_text": response_text, "context": relevant_doc})
    print("img_response is::",img_response)
    format_instructions = parser.get_format_instructions()
    print("format_instructions",format_instructions)
    print("response_text",response_text)

###
    def create_structured_json(img_response):
        result = {}
        for index, img in enumerate(img_response['Image'], start=1):
            print("img",img)
            if img['PageNumber'] is not None:
                # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
                image_key = f"FileName {img['FileName']} PageNumber {img['PageNumber']} ImageNumber {img['ImageNumber']}"
                # Add the image key and description to the result dictionary
                result[f"Image{index}"] = image_key
                result[f"Description{index}"] = img['Description']
            else:
                # Constructing the key format: "file_name_{filename}_page_{page}_image_{image}"
                image_key = f"FileName {img['FileName']} SlideNumber {img['SlideNumber']} ImageNumber {img['ImageNumber']}"
                # Add the image key and description to the result dictionary
                result[f"Image{index}"] = image_key
                result[f"Description{index}"] = img['Description']
        
        return json.dumps(result, indent=4)

    # Using the function to transform the data
    structured_response = create_structured_json(img_response)
    print(structured_response)

    return str(structured_response)
